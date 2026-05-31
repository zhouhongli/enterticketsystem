# -*- coding: utf-8 -*-
"""Generate the information-dense training PPT for the enterprise ticket-system course."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parent.parent / "12-企业售后工单系统构建过程-培训PPT初稿-信息密度增强版.pptx"

FONT = "Microsoft YaHei"
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x11, 0x11, 0x11)
TEXT = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x6F, 0x7C, 0x8E)
LINE = RGBColor(0xDF, 0xE6, 0xEE)
GRID = RGBColor(0xF6, 0xF8, 0xFA)
BLUE = RGBColor(0x14, 0x64, 0xA5)
BLUE_DARK = RGBColor(0x0B, 0x4D, 0x85)
BLUE_PALE = RGBColor(0xEC, 0xF4, 0xFB)
GREEN = RGBColor(0x20, 0xA3, 0x6A)
GREEN_PALE = RGBColor(0xED, 0xF8, 0xF3)
ORANGE = RGBColor(0xF2, 0x8C, 0x28)
ORANGE_PALE = RGBColor(0xFE, 0xF4, 0xE9)

SW = 13.333
SH = 7.5
TOTAL = 40


def I(value):
    return Inches(value)


def east_asian(run):
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    if ea is None:
        ea = OxmlElement("a:ea")
        rpr.append(ea)
    ea.set("typeface", FONT)


def run_style(run, size, color, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    east_asian(run)


def textbox(slide, value, x, y, w, h, size=12, color=TEXT, bold=False,
            align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, spacing=1.03, fit=False):
    shape = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = I(0)
    tf.margin_right = I(0)
    tf.margin_top = I(0)
    tf.margin_bottom = I(0)
    tf.vertical_anchor = valign
    if fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, line_value in enumerate(str(value).split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line_value
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(0)
        for run in p.runs:
            run_style(run, size, color, bold)
    return shape


def rect(slide, x, y, w, h, fill=WHITE, line_color=None, radius=False, width=1):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, I(x), I(y), I(w), I(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(width)
    return shape


def rule(slide, x, y, w, color=BLUE, thickness=0.035):
    return rect(slide, x, y, w, thickness, color)


def line(slide, x1, y1, x2, y2, color=LINE, width=1.2):
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    return shape


def dot(slide, x, y, d, fill, outline=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, I(x), I(y), I(d), I(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if outline is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = outline
        shape.line.width = Pt(1)
    return shape


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide


def chrome(slide, section, number, source=""):
    textbox(slide, section, 0.60, 0.24, 4.8, 0.18, 9.3, BLUE, True)
    rule(slide, 0.60, 0.53, 12.15, LINE, 0.014)
    rule(slide, 0.60, 7.04, 12.15, LINE, 0.014)
    footer = source if source else "从一个想法到可开发方案"
    textbox(slide, footer, 0.60, 7.15, 8.8, 0.20, 8.4, MUTED)
    textbox(slide, f"{number:02d} / {TOTAL:02d}", 11.72, 7.15, 1.04, 0.20, 9.2, MUTED, True, PP_ALIGN.RIGHT)


def body_text(items, numbered=False):
    prefix = [f"{idx + 1}. " for idx in range(len(items))] if numbered else ["• " for _ in items]
    return "\n".join(f"{prefix[idx]}{item}" for idx, item in enumerate(items))


def panel(slide, x, y, w, h, title, items, accent=BLUE, fill=WHITE, numbered=False):
    rect(slide, x, y, w, h, fill, LINE, True)
    rect(slide, x, y, w, 0.40, accent, None)
    textbox(slide, title, x + 0.16, y + 0.12, w - 0.32, 0.15, 10.5, WHITE, True)
    size = 9.1 if len(items) >= 7 else 9.6 if len(items) >= 5 else 10.2
    textbox(slide, body_text(items, numbered), x + 0.18, y + 0.58, w - 0.36, h - 0.74, size, TEXT, False, fit=True)


def tag(slide, value, x, y, w, color=BLUE, fill=BLUE_PALE, size=9.5):
    rect(slide, x, y, w, 0.32, fill, None, True)
    textbox(slide, value, x + 0.04, y + 0.08, w - 0.08, 0.14, size, color, True, PP_ALIGN.CENTER)


def timeline(slide, labels, x, y, w, accent=BLUE):
    gap = 0.10
    node_w = (w - gap * (len(labels) - 1)) / len(labels)
    for idx, label in enumerate(labels):
        nx = x + idx * (node_w + gap)
        fill = GREEN_PALE if idx == len(labels) - 1 else BLUE_PALE
        color = GREEN if idx == len(labels) - 1 else accent
        rect(slide, nx, y, node_w, 0.48, fill, LINE, True)
        textbox(slide, label, nx + 0.04, y + 0.14, node_w - 0.08, 0.14, 8.7, color, True, PP_ALIGN.CENTER)
        if idx < len(labels) - 1:
            line(slide, nx + node_w, y + 0.24, nx + node_w + gap, y + 0.24, MUTED, 1)


def takeaway_for(data):
    return (
        f"{data['left'][0]}；用“{data['mid_title']}”中的项目材料讲例子，"
        f"最后用“{data['right_title']}”作为课堂检查。"
    )


def takeaway(slide, data, y=6.08):
    rect(slide, 0.70, y, 11.92, 0.62, WHITE, LINE, True)
    rule(slide, 0.92, y + 0.19, 0.44, BLUE, 0.04)
    textbox(slide, "本页带走", 1.52, y + 0.20, 1.20, 0.18, 10.3, BLUE, True)
    textbox(slide, takeaway_for(data), 2.78, y + 0.18, 9.36, 0.22, 9.6, TEXT, False, fit=True)


def page_title(slide, title, subtitle):
    textbox(slide, title, 0.62, 0.72, 12.00, 0.40, 23.5, BLACK, True, fit=True)
    textbox(slide, subtitle, 0.64, 1.16, 11.80, 0.22, 10.8, MUTED, False, fit=True)


SLIDES = [
    {
        "section": "开场 / 封面",
        "title": "从一个想法到可开发方案",
        "subtitle": "AI Agent 协作下的软件项目构建全过程，以企业售后工单系统 MVP 为例。",
        "left_title": "课程主线",
        "left": ["不是展示一个完成系统，而是复盘开发前如何形成可执行基线", "从一句想法出发，逐步确认范围、角色、流程、技术和验收", "每个阶段都有可见产物，下一阶段只基于已确认内容推进", "学员可把这套节奏迁移到自己的软件项目"],
        "mid_title": "本案例最终产物",
        "mid": ["产品与技术规范 v1.0", "PRD、SRS、页面与交互说明", "技术设计、数据模型、API 设计", "需求与设计追踪矩阵、测试与验收方案", "MVP 实施计划：尚未进入代码开发"],
        "right_title": "课堂收获",
        "right": ["会问：一个想法如何被澄清", "会写：PRD/SRS/设计/测试如何接续", "会查：如何发现前后文档不一致", "会管：什么时候可以进入编码", "会用 Agent：加快整理，但不放弃人的确认责任"],
        "source": "素材：docs/00-04 已确认项目基线"
    },
    {
        "section": "开场 / 案例定位",
        "title": "一个小 MVP，足够观察完整工程过程",
        "subtitle": "项目不大，但同时包含角色、权限、状态、数据、接口和验收。",
        "left_title": "为什么选它",
        "left": ["售后工单是常见业务，学员容易理解，不需要行业背景", "客户、客服、管理员三类角色天然带出权限边界", "工单状态有明确流转，适合讲状态机和测试", "JSON 存储足够轻，但保留未来数据库演进问题"],
        "mid_title": "项目最小闭环",
        "mid": ["客户注册并登录", "客户创建工单，选择问题分类并填写标题、描述", "管理员创建客服账号并分配工单", "客服处理被分配的工单并公开留言", "工单从待分配推进到处理中、已解决、已关闭"],
        "right_title": "观察点",
        "right": ["范围边界：哪些先做，哪些明确不做", "角色能力：每类用户看到和能操作什么", "后端强校验：不能只靠页面隐藏按钮", "验收证据：不是能演示，而是能验证", "演进设计：MVP 简单但不把未来堵死"],
        "source": "素材：产品与技术规范 v1.0"
    },
    {
        "section": "开场 / 课程边界",
        "title": "本课复盘过程，不推销工具",
        "subtitle": "重点是如何把不确定性一步步压缩，而不是演示某个 AI 工具按钮。",
        "left_title": "本课聚焦",
        "left": ["想法如何被拆成选择题", "MVP 范围如何确认并防止膨胀", "文档之间如何建立承接关系", "验收标准为什么要提前写", "Agent 适合做什么，不适合替人决定什么"],
        "mid_title": "本课不展开",
        "mid": ["不讲 FastAPI 具体代码实现", "不讲前端 CSS 和页面细节编码", "不讲某个 AI 工具安装教程", "不把本案例包装成通用行业方案", "不讨论生产部署、监控和运维体系"],
        "right_title": "讲师使用方式",
        "right": ["先让学员看原始想法，再展示澄清路径", "每章都对照一个真实文档产物", "遇到表格页时强调：这是开发前的共同语言", "避免把 Agent 神化，始终回到基线、确认、验证", "可用自己的项目替换案例里的工单系统"],
        "source": "素材：09/10 培训脚本"
    },
    {
        "section": "开场 / 主线",
        "title": "好的协作，是逐步减少不确定性",
        "subtitle": "每一步不是为了多写文档，而是为了回答下一步必须依赖的问题。",
        "left_title": "过程逻辑",
        "left": ["模糊想法先变成可选择的问题", "用户确认选择后，形成产品与技术基线", "需求文档把业务目标拆成可验收规则", "设计文档把规则落到数据、接口和页面", "测试方案把完成标准提前固定下来"],
        "mid_title": "本项目路径",
        "mid": ["想法：企业售后工单系统", "确认：外部个人客户、三类角色、MVP 范围", "需求：PRD/SRS/页面交互说明", "设计：技术设计/数据模型/API", "验证：追踪矩阵/测试方案/实施计划"],
        "right_title": "可复用公式",
        "right": ["先问服务对象，再问核心闭环", "先定必须做，再定暂不做", "先定角色权限，再谈页面", "先定数据和接口，再安排编码", "先定验收证据，再启动实现"],
        "timeline": ["想法", "澄清", "基线", "设计", "验收", "计划"],
        "source": "素材：总体思路、实施计划"
    },
    {
        "section": "需求澄清 / 起点",
        "title": "起点不是代码，而是一个想法",
        "subtitle": "一句话想法里通常藏着服务对象、流程、权限和非范围。",
        "left_title": "原始表达",
        "left": ["我想做一个企业售后工单系统", "这句话还不能直接进入开发", "它没有说明客户是谁、谁处理、处理到什么程度", "也没有说明哪些常见能力先不做", "如果直接编码，后面会用返工补齐假设"],
        "mid_title": "需要拆出的信息",
        "mid": ["服务对象：外部个人客户还是内部员工", "账号方式：开放注册还是管理员创建", "处理角色：客服是否分组、是否自动分配", "流程边界：提交、处理、解决、关闭如何定义", "技术约束：本地演示还是生产运行"],
        "right_title": "课堂提问",
        "right": ["如果你是产品负责人，会先问哪 5 个问题", "哪些问题必须让业务方选择，不能由 Agent 猜", "哪些能力听起来常见，但不一定属于首期", "一句想法到开发计划之间，至少缺几层产物", "什么情况下可以先做原型，什么情况下不行"],
        "source": "素材：产品与技术规范头脑风暴过程"
    },
    {
        "section": "需求澄清 / 选择题",
        "title": "好问题要让用户做选择，而不是泛泛描述",
        "subtitle": "Agent 的第一价值，是把模糊问题组织成可确认的选项。",
        "left_title": "坏问法",
        "left": ["你想要什么功能", "系统复杂度大概怎样", "需不需要权限", "用什么数据库比较好", "什么时候开始开发"],
        "mid_title": "好问法示例",
        "mid": ["服务对象：个人客户、企业客户、内部员工，先选哪一种", "客户账号：开放注册、邀请码、管理员创建，MVP 采用哪一种", "分配方式：管理员手动分配，还是规则自动分配", "存储方式：先 JSON，本地演示；后续再接数据库", "首期是否包含附件、通知、SLA、报表、搜索"],
        "right_title": "判断标准",
        "right": ["答案能改变范围、数据或权限，就值得问", "问题要有默认推荐，但必须让用户确认", "不要一次问太细，先锁住高影响决策", "确认后的答案要写入基线文档", "未确认的想法只能作为待讨论项"],
        "source": "素材：头脑风暴过程"
    },
    {
        "section": "需求澄清 / 确认",
        "title": "每次确认，都在缩小项目的不确定性",
        "subtitle": "确认不是聊天礼貌，而是后续文档和计划的依据。",
        "left_title": "本项目确认项",
        "left": ["服务对象：外部个人客户", "角色：客户、客服、管理员", "客户账号：开放注册、基础登录/退出", "管理员：维护分类、账号、客户启禁用、分配工单", "客服：查看全部工单，仅处理分配给自己的工单"],
        "mid_title": "确认如何影响设计",
        "mid": ["外部客户意味着必须区分客户门户和内部后台", "开放注册意味着要有客户注册页和密码哈希", "客服只能处理自己的工单意味着后端必须校验负责人", "管理员可处理任意工单意味着需要角色级权限规则", "不做通知意味着验收不检查邮件或短信"],
        "right_title": "落盘要求",
        "right": ["确认项进入产品与技术规范", "未进入 MVP 的内容写成非范围", "技术方向写成确定性约束", "后续 PRD、SRS、设计都引用该基线", "如果调整，必须形成新版本或修订记录"],
        "source": "素材：产品与技术规范 v1.0"
    },
    {
        "section": "MVP 范围 / 必做",
        "title": "MVP 的核心，是最小但完整的业务闭环",
        "subtitle": "不是功能越少越好，而是能证明核心流程成立。",
        "left_title": "必做能力",
        "left": ["客户注册、登录、退出", "客户创建工单并查看自己的工单", "客户在工单内公开留言", "管理员维护问题分类和客服账号", "管理员手动分配或重新分配工单", "客服处理被分配的工单并推进状态"],
        "mid_title": "工单闭环示例",
        "mid": ["客户 A 提交：无法登录系统", "系统创建工单，状态为待分配", "管理员分配给客服 B", "客服 B 留言并改为处理中", "客服 B 解决后改为已解决，再关闭", "客户 A 可查看公开进展，但看不到内部审计字段"],
        "right_title": "检查点",
        "right": ["闭环能否从无数据开始演示", "每个角色是否都有必要入口", "每个状态是否有明确触发者", "客户是否只能访问自己的资源", "禁用账号和退出登录是否会影响访问", "数据重启后是否仍存在"],
        "source": "素材：MVP 范围"
    },
    {
        "section": "MVP 范围 / 暂不做",
        "title": "非范围写清楚，才能保护 MVP",
        "subtitle": "培训中要强调：不做不是遗漏，而是控制首期风险。",
        "left_title": "首期不包含",
        "left": ["附件上传与下载", "邮件、短信、站内通知", "服务级别协议 SLA", "统计报表和管理看板", "高级搜索、多条件筛选", "客户确认关闭或工单重开"],
        "mid_title": "为什么先不做",
        "mid": ["附件会引入文件存储、安全扫描和大小限制", "通知会引入外部服务和失败重试", "SLA 会引入时间统计、超时规则和提醒", "报表会扩大统计口径争议", "高级搜索会推动存储和索引设计提前复杂化"],
        "right_title": "课堂话术",
        "right": ["暂不做不等于永远不做", "把非范围写在文档里，比口头说更能防止范围蔓延", "后续迭代路线可以承接这些能力", "如果业务坚持加入，需要重新评估时间、测试和风险", "Agent 生成建议时也要遵守非范围"],
        "source": "素材：产品与技术规范 v1.0"
    },
    {
        "section": "MVP 范围 / 迭代路线",
        "title": "路线图把“以后再说”变成有顺序的演进",
        "subtitle": "MVP 先验证闭环，后续迭代再扩展效率、统计和正式运行能力。",
        "left_title": "迭代顺序",
        "left": ["MVP：个人客户售后工单闭环和可替换 JSON 仓储", "迭代 1：管理看板，统计总量、分布和客服处理量", "迭代 2：客服效率，多条件筛选、关键词搜索、快捷回复等", "迭代 3：正式运行，数据库、迁移、部署、备份、安全增强"],
        "mid_title": "顺序背后的逻辑",
        "mid": ["没有闭环，统计没有意义", "没有基础数据，搜索和报表无法验证", "小规模试用前，不必提前引入复杂部署", "仓储接口先设计好，后续数据库替换才不会牵动业务逻辑", "账号安全底线从 MVP 就要做，不能推迟"],
        "right_title": "管理价值",
        "right": ["让业务知道哪些诉求会进入后续", "让研发知道当前不必提前实现什么", "让测试知道每期验收边界", "让 Agent 不会把后续能力混入首期", "让变更有谈判依据"],
        "timeline": ["MVP", "看板", "效率", "正式运行"],
        "source": "素材：已确认迭代路线"
    },
    {
        "section": "基线地图 / 文档体系",
        "title": "文档不是堆积，而是逐层回答问题",
        "subtitle": "每份文档都有自己的问题域，混写会让后续实现失去边界。",
        "left_title": "文档链路",
        "left": ["产品与技术规范：确认上位边界", "PRD：说明目标用户、场景和产品目标", "SRS：把产品目标写成可验证的软件规则", "页面交互：说明用户如何操作系统", "技术设计、数据模型、API：说明如何实现这些规则"],
        "mid_title": "继续向下承接",
        "mid": ["追踪矩阵：检查需求、页面、数据和 API 是否一致", "测试与验收方案：定义怎样算做对", "MVP 实施计划：按风险顺序安排开发", "这些文档共同构成开发前基线", "本项目当前尚未开始代码开发"],
        "right_title": "常见误区",
        "right": ["把 PRD 当成所有细节的仓库", "只写页面，不写权限和状态规则", "只写接口，不回溯需求来源", "开发后再补测试用例", "文档确认后不保存版本"],
        "source": "素材：AGENTS.md 必读文档清单"
    },
    {
        "section": "需求基线 / 上位规范",
        "title": "第一份正式产物：产品与技术规范 v1.0",
        "subtitle": "它回答“这个项目到底是什么，第一期做到哪里”。",
        "left_title": "规范回答的问题",
        "left": ["系统服务谁", "MVP 做什么", "MVP 暂不做什么", "角色和状态流程是什么", "采用什么技术方向", "后续按什么路线迭代"],
        "mid_title": "项目中的关键句",
        "mid": ["服务对象：外部个人客户的售后支持", "状态流程：待分配 -> 处理中 -> 已解决 -> 已关闭", "前端：原生 HTML/CSS/JavaScript", "后端：Python FastAPI 提供 REST 接口", "MVP 存储：通过仓储接口访问 JSON 文件"],
        "right_title": "教学重点",
        "right": ["这份文档是后续工作的上位基线", "它不是实现细节，也不是完整需求规格", "但它决定了后续文档不能随意扩大范围", "技术方向必须足够明确，不能停留在建议", "用户确认后才能进入下一层拆解"],
        "source": "素材：产品与技术规范 v1.0"
    },
    {
        "section": "需求基线 / 落盘",
        "title": "口头确认不是项目基线",
        "subtitle": "在 AI 协作中，确认后的内容必须变成可引用、可追踪的文件。",
        "left_title": "只停留在聊天中的问题",
        "left": ["上下文会变长，后续容易遗忘", "同一句话可能被不同人理解为不同范围", "Agent 后续生成内容时可能引用旧假设", "审阅者无法定位正式版本", "项目交接时缺少证据"],
        "mid_title": "本项目做法",
        "mid": ["把确认后的产品与技术基线保存为 v1.0 文档", "把头脑风暴过程单独保存，保留决策背景", "总体思路用于解释目录和开发方法，不替代正式基线", "后续 PRD/SRS/设计只基于正式基线展开", "发现不一致时通过修订关闭"],
        "right_title": "文件示例",
        "right": ["docs/00-项目基线与记录/企业售后工单系统-产品与技术规范-v1.0.md", "docs/00-项目基线与记录/产品与技术规范头脑风暴过程.md", "docs/00-项目基线与记录/总体思路.md"],
        "source": "素材：docs/00 项目基线与记录"
    },
    {
        "section": "需求基线 / PRD",
        "title": "PRD 讲清楚产品目标和用户价值",
        "subtitle": "PRD 不是接口清单，它回答为什么做、给谁用、解决什么场景。",
        "left_title": "PRD 内容",
        "left": ["产品背景和问题描述", "目标用户和典型角色", "核心场景和 MVP 目标", "MVP 非目标", "业务流程和状态规则", "产品验收目标"],
        "mid_title": "项目目标示例",
        "mid": ["G-01：为个人客户提供统一的在线售后问题提交与进展查看入口", "G-03：使负责客服能够基于工单与客户沟通，并推进问题至关闭", "G-05：验证可演进的数据访问架构，为后续真实数据库接入提供基础"],
        "right_title": "讲师提醒",
        "right": ["PRD 里的目标要能继续拆成需求规则", "不要在 PRD 中直接写太多后端实现细节", "非目标必须和 MVP 边界一致", "产品目标后续要能映射到验收用例", "PRD 通过后，SRS 才有稳定输入"],
        "source": "素材：01-产品需求文档-PRD.md"
    },
    {
        "section": "需求基线 / SRS",
        "title": "SRS 把“能做”变成“必须如何做”",
        "subtitle": "它把 PRD 的能力描述拆成角色、条件、字段、异常和状态规则。",
        "left_title": "从 PRD 到 SRS",
        "left": ["PRD：客户应能创建工单", "SRS 要继续回答：谁能创建、创建什么、字段怎么校验", "还要回答：分类是否有效、创建后状态是什么", "还要回答：创建后能否修改、失败时如何提示", "这些规则必须可测试"],
        "mid_title": "需求片段",
        "mid": ["SRS-TKT-001：仅已登录且启用的客户能够创建工单", "SRS-TKT-003：保存前必须再次验证所选分类仍为有效状态", "SRS-TKT-004：创建成功时，状态为待分配，负责人为空", "SRS-SEC：所有角色权限由后端校验", "SRS-DAT：业务逻辑不直接依赖 JSON 文件结构"],
        "right_title": "检查问题",
        "right": ["这条需求有没有明确主体", "这条需求有没有前置条件", "这条需求有没有成功和失败结果", "这条需求能否写成测试用例", "这条需求是否与 PRD 和设计一致"],
        "source": "素材：02-软件需求规格说明书-SRS.md"
    },
    {
        "section": "需求基线 / 页面交互",
        "title": "用户不是操作需求条目，而是操作页面",
        "subtitle": "页面与交互说明把需求变成用户能看到、能点击、能反馈的界面。",
        "left_title": "页面范围",
        "left": ["客户门户：登录、注册、我的工单、新建工单、工单详情", "内部后台：工单列表、工单详情、分类管理、客服账号、客户账号", "每页说明入口、展示字段、操作、状态和异常反馈", "页面规则必须引用 SRS 和 API 设计", "不把客户不应看到的字段放到页面上"],
        "mid_title": "新建工单页示例",
        "mid": ["分类下拉：仅展示启用分类", "标题：必填，1-100 字", "问题描述：必填，1-4000 字", "提交按钮：成功后跳转新建工单详情页", "失败反馈：分类失效、字段为空、未登录或账号禁用"],
        "right_title": "讲师提示",
        "right": ["页面说明不是美术稿，但要足够支持开发", "页面展示字段要和 API 响应字段一致", "权限规则不能只靠页面隐藏按钮", "异常反馈是验收的一部分", "页面文档可作为前后端联调清单"],
        "source": "素材：03-页面与交互说明.md"
    },
    {
        "section": "需求基线 / 追踪关系",
        "title": "重要规则不能只写一次",
        "subtitle": "同一条高风险规则要在需求、页面、API、测试和实施计划中互相照应。",
        "left_title": "追踪规则示例",
        "left": ["客户不展示具体客服账号或负责人名称", "这是数据边界，不只是页面展示问题", "如果 API 返回了字段，前端隐藏也不可靠", "如果测试不覆盖，后续改动容易回归", "如果实施计划不安排复核，风险会被推迟到验收"],
        "mid_title": "映射链条",
        "mid": ["SRS：客户视角不得暴露内部负责人信息", "页面：客户详情页不显示负责人名称", "API：客户详情响应不返回负责人字段", "测试：TC-SEC-004 验证客户响应字段隔离", "实施：M5 安全与字段隔离复核"],
        "right_title": "课堂练习",
        "right": ["任选一条规则，找出它应该出现在哪些文档", "检查是否存在“页面写了，API 没写”的断点", "检查是否存在“需求写了，测试没覆盖”的断点", "把断点记录为一致性审查发现项", "修订后再关闭发现项"],
        "source": "素材：追踪矩阵与一致性审查"
    },
    {
        "section": "系统设计 / 架构",
        "title": "设计不是炫技，而是回应已确认约束",
        "subtitle": "技术设计的核心是让 MVP 可运行、可测试、可替换、可演进。",
        "left_title": "已确认技术方向",
        "left": ["前端：原生 HTML、CSS、JavaScript", "后端：Python FastAPI REST 接口", "存储：通过仓储接口访问 JSON 文件", "安全：密码安全哈希，权限后端校验", "规模：本地演示和小规模试用"],
        "mid_title": "设计回应",
        "mid": ["分层：API 层、服务层、仓储接口、JSON 适配器", "业务规则放在服务层，避免散落在页面或文件读写代码里", "仓储接口隔离数据来源，未来可替换数据库适配器", "认证和授权在后端统一处理", "接口响应按角色裁剪字段"],
        "right_title": "不做的复杂化",
        "right": ["不引入复杂前端框架", "不在 MVP 阶段上真实关系型数据库", "不做微服务拆分", "不做外部消息通知链路", "不把 JSON 文件结构暴露给业务逻辑"],
        "source": "素材：01-技术设计说明书.md"
    },
    {
        "section": "系统设计 / 工程结构",
        "title": "目录结构体现协作边界",
        "subtitle": "好的工程结构让需求、设计、实现和测试能够相互定位。",
        "left_title": "规划原则",
        "left": ["文档先行：docs 目录保存所有已确认基线", "前后端分离但保持同一项目可本地运行", "后端按 API、服务、仓储、模型划分", "测试覆盖关键规则和端到端流程", "数据文件作为 MVP 存储实现，不作为业务模型"],
        "mid_title": "实施计划建议",
        "mid": ["backend/：FastAPI 应用和业务服务", "frontend/：原生页面、样式和交互脚本", "data/：JSON 持久化文件和初始化样例", "tests/：单元、集成、E2E 或验收脚本", "docs/：需求、设计、测试、计划和培训资料"],
        "right_title": "管理价值",
        "right": ["新成员能快速找到依据", "Agent 生成代码时有清晰落点", "测试和实现不会混在一起", "未来数据库替换时影响面可控", "文档与代码演进可以并行审查"],
        "source": "素材：总体思路、MVP实施计划"
    },
    {
        "section": "系统设计 / 数据模型",
        "title": "数据模型把业务规则固化下来",
        "subtitle": "角色、工单、留言、分类、会话和审计都要能支撑已确认流程。",
        "left_title": "核心实体",
        "left": ["User：客户、客服、管理员共用账号实体", "Ticket：工单主体，包含标题、描述、分类、状态、创建人、负责人", "Category：问题分类，由管理员维护启用状态", "TicketMessage：工单公开留言", "Session/Token：登录凭证和失效处理", "AuditLog：关键管理动作和状态变化记录"],
        "mid_title": "字段规则示例",
        "mid": ["password_hash：只保存哈希，不保存明文密码", "ticket.status：只能在待分配、处理中、已解决、已关闭中流转", "assignee_user_id：客户响应中不得暴露", "category.active：创建工单时必须再次校验", "created_at/updated_at：支持审计和后续统计"],
        "right_title": "设计检查",
        "right": ["字段是否支持 SRS 中的所有规则", "字段是否会泄露内部信息", "状态是否能防止非法跳转", "是否保留后续统计所需基础时间", "是否存在 JSON 到数据库的迁移空间"],
        "source": "素材：02-数据模型设计.md"
    },
    {
        "section": "系统设计 / 权限",
        "title": "权限必须由后端强校验",
        "subtitle": "页面隐藏按钮只是体验优化，不是安全边界。",
        "left_title": "角色能力",
        "left": ["客户：创建和查看自己的工单，公开留言", "客服：查看全部工单，按状态筛选，只处理分配给自己的工单", "管理员：维护分类和账号，启禁用客户，分配或重新分配工单", "管理员：可处理任意工单", "禁用客户：不能继续访问受保护能力"],
        "mid_title": "风险场景",
        "mid": ["客户 A 直接访问客户 B 的工单详情 URL", "客服 C 尝试处理分配给客服 D 的工单", "前端被修改后提交非法状态跳转", "客户详情 API 返回 assignee_user_id", "退出登录后仍使用旧 token 调用接口"],
        "right_title": "测试覆盖",
        "right": ["TC-SEC：资源归属校验", "TC-SEC：角色权限矩阵", "TC-AUTH：退出和禁用后的会话失效", "TC-API：客户响应字段隔离", "E2E：越权访问应被拒绝"],
        "source": "素材：SRS、API设计、测试方案"
    },
    {
        "section": "系统设计 / API",
        "title": "API 是前后端协作合同",
        "subtitle": "API 设计要同时承接页面操作、权限规则、数据模型和测试用例。",
        "left_title": "接口分组",
        "left": ["认证：注册、登录、退出、当前用户", "客户：我的工单、新建工单、客户详情、客户留言", "内部：工单列表、分配、状态推进、内部处理", "管理：分类、客服账号、客户账号启禁用", "健康检查和初始化能力"],
        "mid_title": "接口示例",
        "mid": ["POST /api/customer/tickets：客户创建工单", "GET /api/customer/tickets/{id}：客户查看自己的工单", "POST /api/admin/tickets/{id}/assign：管理员分配工单", "PATCH /api/staff/tickets/{id}/status：客服推进状态", "GET /api/staff/tickets?status=processing：内部按状态筛选"],
        "right_title": "合同内容",
        "right": ["请求字段和校验错误", "响应字段按角色裁剪", "状态码和错误码", "权限失败与资源不存在的处理", "接口和测试用例的映射"],
        "source": "素材：03-API接口设计.md"
    },
    {
        "section": "系统设计 / JSON 仓储",
        "title": "JSON 是 MVP 存储，不是业务依赖",
        "subtitle": "轻量存储可以用于演示，但必须通过仓储接口隔离。",
        "left_title": "为什么用 JSON",
        "left": ["本地运行和演示成本低", "不需要提前部署数据库", "便于观察样例数据", "足以支撑客户不超过 100 人、客服不超过 5 人的小规模试用", "适合快速验证业务闭环"],
        "mid_title": "必须避免的做法",
        "mid": ["业务服务直接拼 JSON 文件路径", "页面或 API 层直接读写 JSON", "把 JSON 字段名当成长期数据库字段", "写入时不做原子性保护", "没有初始化、备份或异常恢复思路"],
        "right_title": "设计约束",
        "right": ["定义 Repository 接口", "服务层只依赖接口，不依赖 JSON 实现", "JSON 适配器负责文件读写、锁和错误处理", "后续数据库适配器实现同一接口", "测试要覆盖重启后读取数据"],
        "source": "素材：技术设计、数据模型设计"
    },
    {
        "section": "系统设计 / 状态机",
        "title": "状态流转越简单，越要写清楚",
        "subtitle": "状态规则不写清，工单系统很容易出现跳级、回退和关闭后再改。",
        "left_title": "确认状态",
        "left": ["待分配：客户创建后默认状态，负责人为空", "处理中：管理员分配或客服开始处理后进入", "已解决：客服或管理员标记问题解决", "已关闭：工单闭环结束，MVP 不支持重开", "首期不包含客户确认关闭"],
        "mid_title": "允许与拒绝",
        "mid": ["待分配 -> 处理中：允许", "处理中 -> 已解决：允许", "已解决 -> 已关闭：允许", "待分配 -> 已解决：拒绝", "已关闭 -> 任意状态：拒绝", "已解决 -> 处理中：MVP 拒绝"],
        "right_title": "实现要求",
        "right": ["状态变更由后端服务层统一判断", "非法跳转返回明确错误", "状态变更写入审计记录", "页面按钮只展示可用操作，但不能替代校验", "测试用例覆盖每条高风险转换"],
        "timeline": ["待分配", "处理中", "已解决", "已关闭"],
        "source": "素材：PRD、SRS、测试方案"
    },
    {
        "section": "一致性审查 / 追踪矩阵",
        "title": "审查的目的，是发现文档之间的断点",
        "subtitle": "不是追求文档漂亮，而是确保需求、设计、测试说的是同一件事。",
        "left_title": "审查对象",
        "left": ["PRD 目标是否被 SRS 覆盖", "SRS 规则是否有页面或 API 支撑", "数据模型是否能保存必要字段", "API 是否暴露或隐藏正确字段", "测试方案是否覆盖关键风险"],
        "mid_title": "本项目发现",
        "mid": ["F-01：客户视角负责人展示存在歧义", "F-02：客户建单成功后的跳转不够明确", "F-03：技术设计中仍存在建议性措辞", "三项发现均完成修订并关闭", "一致性审查通过后进入测试与实施计划"],
        "right_title": "课堂操作",
        "right": ["把一条需求横向拉到页面、数据、API、测试", "看到缺口就记录发现项，不要现场脑补", "发现项必须有负责人、修订位置和关闭结论", "通过审查后再把测试方案作为验收依据", "这一步非常适合 Agent 辅助排查"],
        "source": "素材：00-需求与设计追踪矩阵-一致性审查记录.md"
    },
    {
        "section": "一致性审查 / 发现关闭",
        "title": "发现项不关闭，计划就不可靠",
        "subtitle": "开发计划依赖稳定基线；基线有歧义，计划只能排出假确定性。",
        "left_title": "F-01 修订",
        "left": ["问题：客户详情是否展示负责人存在歧义", "风险：暴露内部客服账号或负责人名称", "修订：客户详情不展示具体客服账号或负责人名称", "联动：页面、API、测试均同步确认", "结果：关闭"],
        "mid_title": "F-02 / F-03 修订",
        "mid": ["F-02：创建成功后明确跳转新建工单详情页", "避免前端实现时出现列表页、详情页、弹窗三种不同理解", "F-03：技术方案建议语气改成确定性设计约束", "避免开发时把设计当成可选建议", "三项发现关闭后一致性审查通过"],
        "right_title": "讲师重点",
        "right": ["审查发现不是错误追责，而是开发前排雷", "关闭发现项要改源文档，不只在会议纪要里说明", "涉及多文档的规则要同步修订", "修订后才能进入测试方案确认", "这是 Agent 协作中最容易体现价值的环节"],
        "source": "素材：一致性审查记录"
    },
    {
        "section": "验证前置 / 测试方案",
        "title": "测试不是最后补充，而是提前定义怎样算做对",
        "subtitle": "没有提前定义验收标准，开发完成时就会变成口头争论。",
        "left_title": "测试目标",
        "left": ["验证业务闭环是否完整", "验证权限边界是否可靠", "验证状态流转和数据一致性", "验证凭证安全和密码哈希", "验证 JSON 本地持久化", "验证页面反馈和接口错误处理"],
        "mid_title": "测试层次",
        "mid": ["规则单元：状态、权限、字段校验", "仓储服务：读写、初始化、异常处理", "API 集成：认证、角色、响应字段", "页面联调：主要用户路径", "E2E：从初始化到建单、分配、处理、关闭", "文档/代码审查：安全底线和范围控制"],
        "right_title": "用例示例",
        "right": ["AUTH-004：三类角色登录后进入对应首页", "TKT-005：客户访问他人工单应拒绝", "STS-004：跳级或回退状态应拒绝变更", "DAT-002：应用重启后账号、工单、留言仍存在", "SEC-004：客户响应不包含内部负责人字段"],
        "source": "素材：01-测试与验收方案.md"
    },
    {
        "section": "验证前置 / E2E",
        "title": "验收要覆盖主流程，也要覆盖风险点",
        "subtitle": "E2E 不是点点页面，而是证明关键业务路径和边界同时成立。",
        "left_title": "主流程 E2E",
        "left": ["初始化分类、管理员和客服账号", "客户注册并登录", "客户创建工单并进入详情页", "管理员分配给客服", "客服留言并推进状态", "工单关闭后客户仍可查看历史公开信息"],
        "mid_title": "风险 E2E",
        "mid": ["客户 A 访问客户 B 工单，必须拒绝", "客服 C 修改非本人负责工单，必须拒绝", "禁用客户后旧会话失效", "关闭后的工单不可继续修改", "重启应用后历史数据仍可读取"],
        "right_title": "验收表达",
        "right": ["E2E-01：初始化与建单", "E2E-02：分配、处理、解决、关闭", "E2E-03：越权验证", "E2E-04：禁用客户", "E2E-05：重启读取", "每条 E2E 都要有前置数据、操作步骤和预期结果"],
        "source": "素材：测试与验收方案"
    },
    {
        "section": "验证前置 / 质量门禁",
        "title": "有些缺陷不能靠后续优化处理",
        "subtitle": "安全、权限、数据损坏和闭环不可达属于上线前硬门槛。",
        "left_title": "严重问题",
        "left": ["越权读取或修改他人资源", "明文保存密码或暴露密码哈希", "客户响应泄露内部负责人字段", "状态可非法跳转或关闭后仍可改", "JSON 写入损坏导致数据不可恢复", "主业务闭环无法完成"],
        "mid_title": "验收处理",
        "mid": ["严重：必须修复并复测，否则不得通过", "高：影响安全、流程或数据一致性，必须修复", "中：重要交互反馈缺失，原则上修复", "低：不影响流程的展示问题，可记录后验收", "修复后要回归关联用例"],
        "right_title": "课堂提醒",
        "right": ["把质量门禁写在测试方案里", "不要把安全缺陷归类为体验优化", "不要用“演示时不点那里”规避问题", "Agent 生成代码后也要用门禁检查", "验收不是看页面是否好看，而是看风险是否关闭"],
        "source": "素材：测试与验收方案"
    },
    {
        "section": "实施计划 / 里程碑",
        "title": "计划不是排期表，而是风险处理顺序",
        "subtitle": "MVP 实施计划把已确认基线转成可执行开发路径。",
        "left_title": "M0-M6",
        "left": ["M0：工程骨架和本地启动", "M1：数据、仓储、安全、会话基础", "M2：分类管理、客服账号、客户账号", "M3：客户注册、登录、建单和详情", "M4：内部工单列表、分配、处理、状态流转", "M5：安全、权限、字段隔离与联调", "M6：全量测试、验收和交付材料"],
        "mid_title": "为什么这样排",
        "mid": ["先固化底层规则，避免页面先行造成返工", "先完成账号和权限，再做业务操作", "先完成客户建单，再完成内部处理闭环", "最后集中做安全复核和验收", "每个里程碑都有可验证产物"],
        "right_title": "计划输入",
        "right": ["产品与技术规范", "PRD/SRS/页面交互", "技术设计/数据模型/API", "一致性审查记录", "测试与验收方案", "实施计划确认后再进入工程初始化"],
        "timeline": ["M0", "M1", "M2", "M3", "M4", "M5", "M6"],
        "source": "素材：01-MVP实施计划.md"
    },
    {
        "section": "实施计划 / 拆解逻辑",
        "title": "页面不是最小风险单元",
        "subtitle": "按页面开工看起来快，但容易把权限、状态和数据一致性推到最后。",
        "left_title": "直接按页面的风险",
        "left": ["登录页临时会话", "建单页临时字段校验", "详情页临时权限判断", "管理页临时账号操作", "工单状态在多个页面重复判断", "最后补测试时发现规则分散"],
        "mid_title": "推荐拆解",
        "mid": ["先建立数据模型、仓储接口和基础服务", "再实现认证、会话和权限判定", "再做管理配置和账号能力", "再做客户建单和详情", "再做内部分配、处理和状态流转", "最后做字段隔离、联调和验收"],
        "right_title": "讲师提示",
        "right": ["风险越靠底层，越应提前处理", "页面能跑不等于业务规则正确", "前端开发可以并行，但不能绕过后端规则", "每个阶段要有可运行检查点", "Agent 生成页面前先给它明确接口和权限约束"],
        "source": "素材：MVP实施计划"
    },
    {
        "section": "实施计划 / 测试驱动",
        "title": "高风险规则，适合先写测试",
        "subtitle": "不必所有内容严格 TDD，但权限、状态和数据写入值得测试先行。",
        "left_title": "适合测试先行",
        "left": ["状态转换规则", "角色权限判断", "资源归属校验", "会话失效和禁用客户", "客户响应字段隔离", "JSON 原子写入和重启读取"],
        "mid_title": "测试样例",
        "mid": ["待分配 -> 处理中：允许", "处理中 -> 已解决：允许", "待分配 -> 已解决：拒绝", "已解决 -> 处理中：拒绝", "已关闭 -> 任意状态：拒绝", "客户详情不得包含 assignee_user_id、audit_logs、password_hash"],
        "right_title": "不必强制 TDD",
        "right": ["普通页面排版", "CSS 视觉细节", "一般提示文案微调", "低风险展示字段顺序", "但这些内容仍要经过基本页面检查", "重点是把高风险规则前置"],
        "source": "素材：测试方案、实施计划"
    },
    {
        "section": "实施计划 / 启动检查",
        "title": "不是文档写完了，而是条件满足了",
        "subtitle": "进入编码前要确认基线、设计、测试和计划都已经稳定。",
        "left_title": "已完成条件",
        "left": ["产品与技术规范已确认", "PRD、SRS、页面交互说明已确认", "技术设计、数据模型、API 设计已确认", "需求与设计一致性审查通过", "测试与验收方案已确认"],
        "mid_title": "待确认条件",
        "mid": ["MVP 实施计划 v0.1 待最终确认", "工程目录结构优化已经写入计划", "计划整体确认后方可进入工程初始化", "项目尚未开始代码开发", "进入编码前应再次核对需求与设计文档"],
        "right_title": "启动门槛",
        "right": ["范围没有新增未确认项", "关键规则已有测试策略", "接口与页面没有明显断点", "权限和安全底线已写入计划", "团队知道第一周要完成什么", "用户明确批准进入实现"],
        "source": "素材：AGENTS.md 当前阶段"
    },
    {
        "section": "Agent 协作 / Agent 作用",
        "title": "Agent 的价值，是加快分析和产物衔接",
        "subtitle": "它可以帮助整理、生成、检查和拆解，但不能替业务做最终判断。",
        "left_title": "Agent 适合做",
        "left": ["把模糊想法整理成澄清问题", "根据确认项生成规范文档初稿", "把 PRD 拆成 SRS 和页面交互说明", "根据设计生成 API、数据模型和测试草案", "检查文档之间的不一致", "把验收标准转成实施计划"],
        "mid_title": "本项目中的协作",
        "mid": ["用户提出培训讲师视角", "Agent 先给出页级脚本，再根据反馈增加可见文本", "用户确认制作版脚本后进入 PPT 初稿", "发现页面信息密度不足后，重制增强版", "每一步都保留独立文档和成品"],
        "right_title": "不能交给 Agent 的事",
        "right": ["替业务决定服务对象", "未经确认扩大 MVP 范围", "把建议写成已确认事实", "跳过用户审阅直接进入开发", "忽略安全和权限底线", "把生成速度当成质量保证"],
        "source": "素材：培训资料生成过程"
    },
    {
        "section": "Agent 协作 / 人的责任",
        "title": "Agent 可以推进工作，人必须批准方向",
        "subtitle": "培训中要讲清：AI 协作不是无人负责，而是责任边界更要清楚。",
        "left_title": "人负责",
        "left": ["选择业务场景和目标用户", "确认客户、客服、管理员的权限边界", "确认 MVP 做什么和暂不做什么", "确认技术约束是否符合团队能力", "批准基线文档进入下一阶段", "判断质量是否满足真实授课或开发需要"],
        "mid_title": "Agent 协助",
        "mid": ["提出选项和风险", "整理头脑风暴过程", "生成结构化文档", "根据反馈重写和扩写", "用矩阵检查一致性", "自动生成 PPT 初稿和渲染抽查材料"],
        "right_title": "协作原则",
        "right": ["Agent 先给可审阅产物", "人基于业务判断做确认", "确认后写入正式文件", "发现问题时修订源文档或源脚本", "不要让聊天记录替代版本化文档", "最终质量由人验收"],
        "source": "素材：项目协作记录"
    },
    {
        "section": "Agent 协作 / 过程纪律",
        "title": "Agent 协作也需要过程纪律",
        "subtitle": "没有纪律，AI 只会更快地产生更多不一致内容。",
        "left_title": "本项目纪律",
        "left": ["确认后要落盘", "建议措辞要定稿", "发现项要关闭", "非范围不要偷做", "测试标准不能后补", "培训资料独立于 MVP 基线"],
        "mid_title": "纠偏示例",
        "mid": ["旧版 PPT 页级脚本信息深度不足，被保留为历史材料", "四页视觉样稿信息量不足，保留供对照", "新版页级脚本和讲师稿重新展开三种模式和案例", "本次企业工单系统 PPT 初稿字数偏少，进入信息密度增强版", "每次反馈都转化为具体新产物"],
        "right_title": "讲师可强调",
        "right": ["AI 生成不是一次到位", "反馈要具体：信息量、示例、可见文本、页面充实度", "保留旧版本有助于讲解迭代过程", "修订时先明确目标，再改脚本和成品", "生成后必须渲染检查，不只看文件存在"],
        "source": "素材：docs/05 培训资料过程"
    },
    {
        "section": "Agent 协作 / 可复用模式",
        "title": "用同一套节奏处理下一个项目",
        "subtitle": "这套方法不依赖工单系统，适合多数从零开始的软件项目。",
        "left_title": "八步节奏",
        "left": ["问清场景", "确认 MVP", "写清非范围", "固化基线", "拆需求和页面", "展开数据、接口和技术设计", "先定测试与验收", "再排实施计划"],
        "mid_title": "迁移到其他项目",
        "mid": ["内部知识库：知识条目、搜索、审核、权限、版本", "培训报名系统：课程、学员、报名、支付、签到", "设备维修系统：报修、派工、处理、验收、统计", "每个项目都要先确认服务对象和最小闭环", "不要一开始就讨论全部高级功能"],
        "right_title": "课堂练习提示",
        "right": ["让学员拿自己的项目写一句原始想法", "设计 5 个高影响澄清问题", "写 3 项 MVP 必做和 5 项暂不做", "画出一条主业务闭环", "定义一条 E2E 验收路径", "最后再讨论技术选型"],
        "source": "素材：培训大纲与页级脚本"
    },
    {
        "section": "收束 / 总结",
        "title": "规格清楚，Agent 才能跑得稳",
        "subtitle": "速度来自清晰的约束，而不是跳过思考。",
        "left_title": "核心结论",
        "left": ["范围有边界，开发才不会无限膨胀", "角色有权限，系统才不会靠页面假安全", "状态有规则，流程才不会随意跳转", "数据有模型，接口才不会各说各话", "测试有证据，验收才不会变成感觉判断"],
        "mid_title": "本项目证明",
        "mid": ["从一句想法形成产品与技术规范", "从规范展开 PRD、SRS、页面、设计和 API", "通过追踪矩阵发现并关闭三项一致性问题", "测试与验收方案在开发前确认", "实施计划等待最终确认后才进入编码"],
        "right_title": "留给学员",
        "right": ["不要把 AI 当成需求来源", "不要把文档当成形式主义", "不要把测试推到最后", "不要让后续迭代混入 MVP", "最重要的是：每一步都要有人确认"],
        "source": "素材：全项目过程复盘"
    },
    {
        "section": "收束 / 课堂练习",
        "title": "练习：把一句想法变成 MVP 边界",
        "subtitle": "让学员现场体会：真正困难的是把范围说清楚。",
        "left_title": "练习步骤",
        "left": ["写一句原始想法", "设计 5 个澄清问题", "写 3 项 MVP 必做", "写 5 项 MVP 暂不做", "定义 1 条 E2E 验收路径", "指出至少 2 个权限或数据风险"],
        "mid_title": "示例题",
        "mid": ["原始想法：做一个内部知识库系统", "目标用户：内部员工、知识管理员", "MVP 必做：创建知识条目、浏览和搜索、管理员审核", "MVP 暂不做：全文智能问答、复杂权限组、外部分享、积分排行、移动端 App", "E2E：员工提交知识，管理员审核通过，其他员工搜索并查看"],
        "right_title": "评价标准",
        "right": ["是否有明确服务对象", "MVP 是否形成闭环", "非范围是否足够具体", "验收路径是否可操作", "是否识别权限和数据边界", "是否能进入下一步 PRD/SRS 拆解"],
        "source": "素材：培训活动设计"
    },
    {
        "section": "收束 / 讨论",
        "title": "讨论：你的团队最容易跳过哪一步？",
        "subtitle": "用最后一页把课程拉回学员自己的项目现场。",
        "left_title": "常见跳过项",
        "left": ["需求澄清：一开始默认大家理解一致", "非范围确认：只写做什么，不写暂不做", "设计追踪：需求、页面、数据、API 缺少映射", "测试前置：代码后才讨论完成标准", "变更确认：开发中悄悄扩大范围"],
        "mid_title": "讨论问题",
        "mid": ["你所在团队最常见的返工来自哪里", "哪些规则经常靠口头约定", "哪些需求最容易被页面原型掩盖", "AI 介入后，是更清楚了还是更容易跳步骤", "你准备把哪一页方法带回团队试用"],
        "right_title": "收束句",
        "right": ["Agent 加快执行，人负责判断", "基线越清楚，协作越可靠", "测试越前置，验收越少争论", "非范围越明确，MVP 越容易完成", "从下一次需求讨论开始，就把确认写下来"],
        "source": "Q&A"
    },
]


def build():
    prs = Presentation()
    prs.slide_width = I(SW)
    prs.slide_height = I(SH)

    for idx, data in enumerate(SLIDES, start=1):
        slide = new_slide(prs)
        chrome(slide, data["section"], idx, data.get("source", ""))

        if idx == 1:
            rule(slide, 0.72, 0.84, 0.88, BLUE, 0.055)
            textbox(slide, data["title"], 0.72, 1.12, 7.8, 0.58, 32, BLACK, True)
            textbox(slide, "信息密度增强版", 8.92, 1.18, 2.60, 0.30, 15, GREEN, True, PP_ALIGN.CENTER)
            textbox(slide, data["subtitle"], 0.76, 1.92, 8.3, 0.30, 15, BLUE_DARK, True)
            panel(slide, 0.76, 2.58, 3.80, 3.38, data["left_title"], data["left"], BLUE, BLUE_PALE)
            panel(slide, 4.78, 2.58, 4.05, 3.38, data["mid_title"], data["mid"], GREEN, GREEN_PALE)
            panel(slide, 9.05, 2.58, 3.55, 3.38, data["right_title"], data["right"], ORANGE, ORANGE_PALE)
            timeline(slide, ["想法", "澄清", "基线", "设计", "验收", "计划"], 1.10, 6.30, 10.90)
            continue

        page_title(slide, data["title"], data["subtitle"])
        panel(slide, 0.70, 1.68, 3.65, 4.18, data["left_title"], data["left"], BLUE, BLUE_PALE)
        panel(slide, 4.56, 1.68, 4.25, 4.18, data["mid_title"], data["mid"], GREEN, GREEN_PALE)
        panel(slide, 9.02, 1.68, 3.60, 4.18, data["right_title"], data["right"], ORANGE, ORANGE_PALE)
        if "timeline" in data:
            timeline(slide, data["timeline"], 1.02, 6.44, 11.20)
        else:
            takeaway(slide, data)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(str(OUT))
    print(f"slides={len(prs.slides)}")
    return OUT


if __name__ == "__main__":
    build()
