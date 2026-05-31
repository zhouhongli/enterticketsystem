# -*- coding: utf-8 -*-
"""Generate four redesigned visual sample slides for review."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parent.parent / "05-新版视觉样稿-四页.pptx"

FONT = "Microsoft YaHei"
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x0E, 0x18, 0x25)
TEXT = RGBColor(0x34, 0x40, 0x52)
MUTED = RGBColor(0x6F, 0x7C, 0x8E)
LINE = RGBColor(0xE4, 0xEA, 0xF0)
BLUE = RGBColor(0x07, 0x56, 0xC7)
BLUE_DARK = RGBColor(0x04, 0x3D, 0x93)
BLUE_PALE = RGBColor(0xF1, 0xF6, 0xFC)
GREEN = RGBColor(0x0C, 0xA4, 0x69)
GREEN_PALE = RGBColor(0xED, 0xF8, 0xF3)
ORANGE = RGBColor(0xF3, 0x83, 0x21)
ORANGE_PALE = RGBColor(0xFE, 0xF4, 0xE9)
GRID = RGBColor(0xF2, 0xF5, 0xF8)

SW = 13.333
SH = 7.5


def I(n):
    return Inches(n)


def east_asian(run):
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    if ea is None:
        ea = OxmlElement("a:ea")
        rpr.append(ea)
    ea.set("typeface", FONT)


def run_style(run, size, color, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    east_asian(run)


def text(slide, value, x, y, w, h, size=18, color=TEXT, bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE, spacing=1.05):
    shape = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = I(0.0)
    tf.margin_right = I(0.0)
    tf.margin_top = I(0.0)
    tf.margin_bottom = I(0.0)
    tf.vertical_anchor = valign
    for idx, line in enumerate(value.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(0)
        for r in p.runs:
            run_style(r, size, color, bold)
    return shape


def rect(slide, x, y, w, h, fill=WHITE, line=None, radius=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(kind, I(x), I(y), I(w), I(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    return shp


def rule(slide, x, y, w, color=BLUE, thickness=0.04):
    return rect(slide, x, y, w, thickness, color)


def line(slide, x1, y1, x2, y2, color=LINE, width=1.25):
    shp = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    shp.line.color.rgb = color
    shp.line.width = Pt(width)
    return shp


def dot(slide, x, y, d, fill, outline=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, I(x), I(y), I(d), I(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if outline:
        shp.line.color.rgb = outline
        shp.line.width = Pt(1.25)
    else:
        shp.line.fill.background()
    return shp


def tag(slide, value, x, y, w, color=BLUE, fill=BLUE_PALE):
    rect(slide, x, y, w, 0.34, fill, None, True)
    text(slide, value, x, y + 0.065, w, 0.22, 11, color, True, PP_ALIGN.CENTER)


def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    return s


def page_chrome(s, label, number):
    text(s, label, 0.62, 0.34, 2.3, 0.19, 10, BLUE, True)
    rule(s, 0.62, 0.64, 12.07, LINE, 0.018)
    rule(s, 0.62, 7.05, 12.07, LINE, 0.018)
    text(s, "传统软件开发模式与 AI Agent 协作开发", 0.62, 7.15, 4.20, 0.18, 9, MUTED)
    text(s, f"S{number:02d}  /  04", 11.84, 7.14, 0.84, 0.18, 10, MUTED, True, PP_ALIGN.RIGHT)


def numbered_dot(s, value, x, y, accent=BLUE, d=0.30):
    dot(s, x, y, d, WHITE, accent)
    text(s, value, x, y + 0.065, d, 0.18, 10, accent, True, PP_ALIGN.CENTER)


def sample_01(prs):
    s = slide(prs)
    page_chrome(s, "技术管理培训  /  视觉样稿", 1)
    # Minimal typographic hero.
    text(s, "软件开发模式，", 0.62, 1.14, 6.40, 0.62, 40, BLACK, True)
    text(s, "正在进入", 0.62, 1.86, 3.20, 0.62, 40, BLACK, True)
    text(s, "AI Agent 时代", 3.82, 1.86, 5.54, 0.62, 40, BLUE, True)
    text(s, "传统软件开发模式与 AI Agent 时代的软件开发模式", 0.64, 2.78, 7.40, 0.34, 18, TEXT)
    text(s, "面向产品、项目与研发管理人员", 0.64, 3.31, 5.10, 0.30, 16, MUTED)
    # Geometric right-hand structure: no images, only a system line.
    text(s, "FROM PROCESS  TO  COLLABORATION", 8.14, 1.29, 4.53, 0.24, 10, MUTED, True)
    line(s, 8.14, 2.01, 12.41, 2.01, BLUE, 2)
    stages = [
        ("01", "阶段控制", "需求与责任", BLUE),
        ("02", "持续反馈", "迭代与交付", BLUE),
        ("03", "Agent 协作", "执行与验证", GREEN),
    ]
    for i, (no, head, sub, accent) in enumerate(stages):
        yy = 2.36 + i * 1.13
        numbered_dot(s, no, 8.16, yy + 0.06, accent, 0.34)
        line(s, 8.58, yy + 0.23, 9.10, yy + 0.23, accent, 1.3)
        text(s, head, 9.29, yy, 2.65, 0.29, 18, BLACK, True)
        text(s, sub, 9.29, yy + 0.39, 2.65, 0.25, 14, MUTED)
    rule(s, 0.64, 5.33, 0.76, BLUE, 0.055)
    text(s, "工程纪律被保留，协作颗粒度被重塑", 0.64, 5.59, 7.14, 0.36, 21, BLUE_DARK, True)
    tag(s, "单位 / 讲师 / 日期 待补充", 0.64, 6.37, 2.76, BLUE, BLUE_PALE)
    return s


def sample_02(prs):
    s = slide(prs)
    page_chrome(s, "模式演进  /  总图", 2)
    text(s, "软件开发模式，并非推倒重来", 0.62, 0.92, 8.80, 0.48, 36, BLACK, True)
    text(s, "每一次演进，都在回答新的协作瓶颈", 0.64, 1.54, 8.60, 0.30, 16, MUTED)
    # timeline backbone
    line(s, 0.72, 3.15, 12.18, 3.15, BLUE, 2.25)
    eras = [
        ("01", "瀑布 / V 模型", "控制复杂度", "基线 · 评审 · 责任", 0.84, BLUE),
        ("02", "敏捷 / DevOps", "缩短反馈周期", "迭代 · 流水线 · 反馈", 4.50, BLUE),
        ("03", "AI Agent 协作", "细化可委派工作", "上下文 · 执行 · 验证", 8.34, GREEN),
    ]
    for no, head, issue, retains, x, accent in eras:
        dot(s, x, 2.90, 0.50, WHITE, accent)
        text(s, no, x, 3.065, 0.50, 0.18, 11, accent, True, PP_ALIGN.CENTER)
        text(s, head, x, 2.11, 3.25, 0.34, 20, BLACK, True)
        text(s, issue, x, 3.76, 3.30, 0.31, 18, accent, True)
        text(s, retains, x, 4.27, 3.35, 0.29, 15, TEXT)
        rule(s, x, 4.81, 0.62, accent, 0.045)
    # lower message uses less boxing, more hierarchy.
    line(s, 0.72, 5.56, 12.18, 5.56, LINE, 1)
    text(s, "保留", 0.78, 5.88, 0.60, 0.24, 12, MUTED, True)
    text(s, "工程基线", 1.60, 5.83, 1.56, 0.32, 18, BLUE_DARK, True)
    text(s, "+", 3.44, 5.84, 0.25, 0.30, 18, MUTED, True)
    text(s, "持续反馈", 3.98, 5.83, 1.56, 0.32, 18, BLUE_DARK, True)
    text(s, "+", 5.86, 5.84, 0.25, 0.30, 18, MUTED, True)
    text(s, "Agent 加速", 6.39, 5.83, 1.78, 0.32, 18, GREEN, True)
    rect(s, 9.13, 5.69, 3.05, 0.65, GREEN_PALE, None, True)
    text(s, "形成可靠交付", 9.13, 5.88, 3.05, 0.28, 18, GREEN, True, PP_ALIGN.CENTER)
    return s


def sample_03(prs):
    s = slide(prs)
    page_chrome(s, "全生命周期  /  Agent 参与", 3)
    text(s, "Agent 进入的不只是编码环节", 0.62, 0.92, 8.80, 0.50, 36, BLACK, True)
    text(s, "执行加速贯穿全流程，关键判断仍由人把关", 0.64, 1.54, 9.20, 0.30, 16, MUTED)
    # Two-lane map
    text(s, "人的门禁", 0.66, 2.26, 1.00, 0.24, 13, BLUE_DARK, True)
    text(s, "Agent 参与", 0.66, 4.30, 1.03, 0.24, 13, GREEN, True)
    line(s, 1.96, 2.43, 12.28, 2.43, BLUE, 1.6)
    line(s, 1.96, 4.48, 12.28, 4.48, GREEN, 1.6)
    phases = [
        ("需求", "目标确认", "整理分歧"),
        ("设计", "架构批准", "比较方案"),
        ("编码", "范围审查", "执行修改"),
        ("测试", "通过判定", "运行验证"),
        ("发布", "上线批准", "核对清单"),
        ("运维", "处置决策", "归纳日志"),
        ("反馈", "优先决策", "汇总建议"),
    ]
    for i, (phase, human, agent) in enumerate(phases):
        x = 2.10 + i * 1.47
        dot(s, x, 2.30, 0.26, WHITE, BLUE)
        dot(s, x, 4.35, 0.26, WHITE, GREEN)
        line(s, x + 0.13, 2.57, x + 0.13, 4.34, LINE, 0.9)
        text(s, phase, x - 0.24, 3.29, 0.74, 0.27, 16, BLACK, True, PP_ALIGN.CENTER)
        text(s, human, x - 0.44, 1.87, 1.13, 0.23, 11, BLUE_DARK, True, PP_ALIGN.CENTER)
        text(s, agent, x - 0.44, 4.78, 1.13, 0.23, 11, GREEN, True, PP_ALIGN.CENTER)
    rule(s, 0.67, 5.80, 0.72, BLUE, 0.05)
    text(s, "加速执行，不等于让渡决策", 0.67, 6.04, 5.52, 0.36, 22, BLUE_DARK, True)
    tag(s, "目标 · 权限 · 验证 · 责任", 9.08, 6.04, 3.12, BLUE, BLUE_PALE)
    return s


def sample_04(prs):
    s = slide(prs)
    page_chrome(s, "治理结论  /  五项原则", 4)
    text(s, "速度提高之后，", 0.62, 0.98, 5.00, 0.55, 40, BLACK, True)
    text(s, "控制原则更重要", 0.62, 1.58, 5.60, 0.55, 40, BLUE, True)
    text(s, "Agent 能加速正确工作，也能加速错误落地", 0.64, 2.40, 7.80, 0.31, 17, TEXT)
    # Five governance principles as a single system, not cards.
    items = [
        ("01", "目标明确", "明确成功标准", BLUE),
        ("02", "上下文充分", "提供真实约束", BLUE),
        ("03", "授权有界", "控制操作范围", ORANGE),
        ("04", "验证有据", "保留测试证据", GREEN),
        ("05", "责任有人", "批准最终结果", BLUE),
    ]
    line(s, 0.87, 4.05, 12.04, 4.05, LINE, 1.5)
    for i, (no, label, note, accent) in enumerate(items):
        x = 0.90 + i * 2.32
        dot(s, x, 3.88, 0.35, accent)
        text(s, no, x, 3.99, 0.35, 0.15, 9, WHITE, True, PP_ALIGN.CENTER)
        text(s, label, x, 4.50, 1.78, 0.30, 18, accent, True)
        text(s, note, x, 4.94, 1.82, 0.26, 13, MUTED)
    rect(s, 0.65, 6.02, 12.02, 0.60, BLUE_DARK, None, False)
    text(s, "可控协作，才是可靠加速", 0.65, 6.17, 12.02, 0.28, 21, WHITE, True, PP_ALIGN.CENTER)
    return s


def build():
    prs = Presentation()
    prs.slide_width = I(SW)
    prs.slide_height = I(SH)
    prs.core_properties.title = "新版视觉样稿：传统软件开发模式与 AI Agent 时代的软件开发模式"
    prs.core_properties.author = "Codex"
    sample_01(prs)
    sample_02(prs)
    sample_03(prs)
    sample_04(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(str(OUT))
    return OUT


if __name__ == "__main__":
    build()

