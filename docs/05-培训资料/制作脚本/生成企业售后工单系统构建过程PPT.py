# -*- coding: utf-8 -*-
"""Generate the first draft PPT for the enterprise ticket-system build-process course."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parent.parent / "11-企业售后工单系统构建过程-培训PPT初稿.pptx"

FONT = "Microsoft YaHei"
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x11, 0x11, 0x11)
TEXT = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x6F, 0x7C, 0x8E)
LINE = RGBColor(0xE4, 0xEA, 0xF0)
GRID = RGBColor(0xF6, 0xF8, 0xFA)
BLUE = RGBColor(0x14, 0x64, 0xA5)
BLUE_DARK = RGBColor(0x0B, 0x4D, 0x85)
BLUE_PALE = RGBColor(0xEC, 0xF4, 0xFB)
GREEN = RGBColor(0x20, 0xA3, 0x6A)
GREEN_PALE = RGBColor(0xED, 0xF8, 0xF3)
ORANGE = RGBColor(0xF2, 0x8C, 0x28)
ORANGE_PALE = RGBColor(0xFE, 0xF4, 0xE9)
RED_PALE = RGBColor(0xFE, 0xEE, 0xED)

SW = 13.333
SH = 7.5
TOTAL = 40


def I(n):
    return Inches(n)


def east_asian(run):
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    if ea is None:
        ea = OxmlElement("a:ea")
        rpr.append(ea)
    ea.set("typeface", FONT)


def run_style(run, size, color, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    east_asian(run)


def textbox(slide, value, x, y, w, h, size=16, color=TEXT, bold=False,
            align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, spacing=1.04):
    shape = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = I(0)
    tf.margin_right = I(0)
    tf.margin_top = I(0)
    tf.margin_bottom = I(0)
    tf.vertical_anchor = valign
    for idx, line_value in enumerate(value.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line_value
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(0)
        for run in p.runs:
            run_style(run, size, color, bold)
    return shape


def rect(slide, x, y, w, h, fill=WHITE, line_color=None, radius=False, width=1):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, I(x), I(y), I(w), I(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(width)
    return shape


def rule(slide, x, y, w, color=BLUE, thickness=0.035):
    return rect(slide, x, y, w, thickness, color)


def line(slide, x1, y1, x2, y2, color=LINE, width=1.2):
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    return shape


def dot(slide, x, y, d, fill, outline=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, I(x), I(y), I(d), I(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if outline is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = outline
        shape.line.width = Pt(1.1)
    return shape


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide


def chrome(slide, section, number):
    textbox(slide, section, 0.60, 0.28, 4.2, 0.20, 9.5, BLUE, True)
    rule(slide, 0.60, 0.56, 12.15, LINE, 0.015)
    rule(slide, 0.60, 7.04, 12.15, LINE, 0.015)
    textbox(slide, "从一个想法到可开发方案", 0.60, 7.15, 4.0, 0.20, 9, MUTED)
    textbox(slide, f"{number:02d} / {TOTAL:02d}", 11.76, 7.15, 1.0, 0.20, 9.5, MUTED, True, PP_ALIGN.RIGHT)


def page(prs, section, number, title, subtitle=None, key=False):
    slide = new_slide(prs)
    chrome(slide, section, number)
    size = 34 if key else 27
    textbox(slide, title, 0.62, 0.80, 11.8, 0.48, size, BLACK, True)
    if subtitle:
        textbox(slide, subtitle, 0.64, 1.36 if not key else 1.50, 11.2, 0.24, 14.5, MUTED)
    return slide


def bullet(slide, value, x, y, w, size=15.5, accent=BLUE, color=TEXT):
    dot(slide, x, y + 0.10, 0.07, accent)
    textbox(slide, value, x + 0.16, y, w - 0.16, 0.30, size, color)


def tag(slide, value, x, y, w, color=BLUE, fill=BLUE_PALE, size=10.5):
    rect(slide, x, y, w, 0.32, fill, None, True)
    textbox(slide, value, x + 0.03, y + 0.07, w - 0.06, 0.16, size, color, True, PP_ALIGN.CENTER)


def card(slide, x, y, w, h, header, lines, accent=BLUE, fill=WHITE, body_size=14.5):
    rect(slide, x, y, w, h, fill, LINE, True)
    rule(slide, x + 0.18, y + 0.18, 0.46, accent, 0.045)
    textbox(slide, header, x + 0.18, y + 0.35, w - 0.36, 0.28, 17, BLACK, True)
    for idx, item in enumerate(lines):
        bullet(slide, item, x + 0.22, y + 0.78 + idx * 0.34, w - 0.38, body_size, accent)


def mini_card(slide, x, y, w, h, header, body, accent=BLUE, fill=WHITE, head_size=15, body_size=12.2):
    rect(slide, x, y, w, h, fill, LINE, True)
    rule(slide, x + 0.14, y + 0.14, 0.38, accent, 0.04)
    if body and h <= 0.95:
        textbox(slide, header, x + 0.14, y + 0.23, w - 0.28, 0.20, head_size, BLACK, True)
        textbox(slide, body, x + 0.14, y + 0.48, w - 0.28, max(0.18, h - 0.54), body_size, TEXT)
    else:
        textbox(slide, header, x + 0.14, y + 0.31, w - 0.28, 0.23, head_size, BLACK, True)
        if body:
            textbox(slide, body, x + 0.14, y + 0.66, w - 0.28, h - 0.78, body_size, TEXT)


def code_box(slide, value, x, y, w, h, size=10.5):
    rect(slide, x, y, w, h, GRID, LINE, True)
    textbox(slide, value, x + 0.16, y + 0.14, w - 0.32, h - 0.26, size, TEXT)


def table(slide, x, y, widths, headers, rows, row_h=0.45, size=11.2, accent=BLUE):
    total_w = sum(widths)
    rect(slide, x, y, total_w, row_h, accent, None)
    cx = x
    for idx, h in enumerate(headers):
        textbox(slide, h, cx + 0.06, y + 0.13, widths[idx] - 0.12, row_h - 0.15, size, WHITE, True, PP_ALIGN.CENTER)
        cx += widths[idx]
    cy = y + row_h
    for r_idx, row in enumerate(rows):
        fill = WHITE if r_idx % 2 == 0 else GRID
        rect(slide, x, cy, total_w, row_h, fill, LINE)
        cx = x
        for c_idx, cell in enumerate(row):
            textbox(slide, str(cell), cx + 0.07, cy + 0.11, widths[c_idx] - 0.14, row_h - 0.12, size, TEXT)
            cx += widths[c_idx]
        cy += row_h


def timeline(slide, labels, x, y, width, node_h=0.56, accent=BLUE, size=12.2):
    count = len(labels)
    gap = 0.14
    node_w = (width - gap * (count - 1)) / count
    for idx, label in enumerate(labels):
        nx = x + idx * (node_w + gap)
        rect(slide, nx, y, node_w, node_h, BLUE_PALE, None, True)
        textbox(slide, label, nx + 0.03, y + 0.12, node_w - 0.06, node_h - 0.18, size, accent, True, PP_ALIGN.CENTER)
        if idx < count - 1:
            textbox(slide, ">", nx + node_w - 0.01, y + 0.17, gap + 0.02, 0.22, 12, MUTED, True, PP_ALIGN.CENTER)


def make_prs():
    prs = Presentation()
    prs.slide_width = I(SW)
    prs.slide_height = I(SH)
    return prs


def build():
    prs = make_prs()

    # P01
    s = new_slide(prs)
    rect(s, 0, 0, SW, SH, WHITE)
    rule(s, 0.76, 0.74, 0.82, BLUE, 0.055)
    textbox(s, "从一个想法到可开发方案", 0.76, 1.12, 8.8, 0.62, 36, BLACK, True)
    textbox(s, "AI Agent 协作下的软件项目构建全过程", 0.80, 1.92, 8.2, 0.34, 20, BLUE_DARK, True)
    textbox(s, "以企业售后工单系统 MVP 为例", 0.80, 2.40, 5.5, 0.28, 16, MUTED)
    timeline(s, ["想法澄清", "需求基线", "设计基线", "测试验收", "实施计划"], 0.86, 4.60, 10.8, 0.60)
    textbox(s, "产品经理 / 项目经理 / 研发管理人员 / 业务分析人员", 0.84, 6.56, 7.4, 0.22, 11, MUTED)
    textbox(s, "01 / 40", 11.74, 6.56, 0.84, 0.22, 11, MUTED, True, PP_ALIGN.RIGHT)

    # P02
    s = page(prs, "开场 / 案例定位", 2, "一个小 MVP，足够观察完整工程过程", "它不大，但包含真实软件项目常见的风险点。")
    rect(s, 0.78, 2.08, 3.55, 3.46, BLUE_DARK, None, True)
    textbox(s, "企业售后\n工单系统\nMVP", 1.08, 2.70, 2.95, 1.24, 29, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(s, "从想法到开发前基线", 1.00, 4.60, 3.10, 0.27, 14, BLUE_PALE, True, PP_ALIGN.CENTER)
    obs = [
        ("三类角色", "客户、客服、管理员"),
        ("权限边界", "客户本人 / 客服负责人 / 管理员全部"),
        ("状态流转", "待分配到已关闭"),
        ("数据持久化", "JSON 与未来数据库"),
        ("验收证据", "权限、安全、状态、持久化"),
    ]
    for i, (h, b) in enumerate(obs):
        y = 1.86 + i * 0.72
        tag(s, f"{i+1}", 4.85, y + 0.03, 0.42, GREEN if i == 4 else BLUE, GREEN_PALE if i == 4 else BLUE_PALE)
        textbox(s, h, 5.42, y, 2.0, 0.25, 17, BLACK, True)
        textbox(s, b, 7.36, y + 0.03, 4.52, 0.22, 13, TEXT)

    # P03
    s = page(prs, "开场 / 课程边界", 3, "本课复盘过程，不推销工具")
    card(s, 0.80, 1.78, 5.58, 4.45, "本课聚焦", ["想法如何被问清", "MVP 范围如何确认", "文档如何衔接", "验收如何前置", "Agent 如何协作", "人如何负责判断"], BLUE, BLUE_PALE, 15)
    card(s, 6.78, 1.78, 5.58, 4.45, "本课不展开", ["FastAPI 代码实现", "前端编码细节", "AI 工具安装教程", "产品售卖", "生产部署运维"], ORANGE, ORANGE_PALE, 15)

    # P04
    s = page(prs, "开场 / 主线", 4, "好的协作，是逐步减少不确定性")
    labels = ["模糊想法", "选择题澄清", "产品基线", "需求与设计", "测试验收", "实施计划"]
    subs = ["一句话", "服务谁/做什么/不做什么", "规范 v1.0", "PRD/SRS/设计", "矩阵/测试方案", "M0-M6"]
    line(s, 1.0, 3.25, 12.0, 3.25, BLUE, 2.0)
    for i, lab in enumerate(labels):
        x = 0.82 + i * 2.03
        dot(s, x + 0.44, 3.05, 0.42, WHITE, BLUE if i < 5 else GREEN)
        textbox(s, str(i + 1), x + 0.44, 3.16, 0.42, 0.14, 9, BLUE if i < 5 else GREEN, True, PP_ALIGN.CENTER)
        textbox(s, lab, x, 2.30, 1.28, 0.24, 15, BLACK, True, PP_ALIGN.CENTER)
        textbox(s, subs[i], x - 0.10, 3.78, 1.48, 0.38, 11.5, MUTED, False, PP_ALIGN.CENTER)
    textbox(s, "前期工作不是“多写文档”，而是持续回答关键问题。", 1.08, 5.45, 11.0, 0.34, 19, BLUE_DARK, True, PP_ALIGN.CENTER)

    # P05
    s = page(prs, "需求澄清 / 起点", 5, "起点不是代码，而是一个想法")
    rect(s, 0.80, 1.86, 4.85, 3.80, GRID, LINE, True)
    textbox(s, "项目初始状态", 1.05, 2.15, 2.5, 0.28, 18, BLACK, True)
    for i, item in enumerate(["工作目录为空", "没有已有代码", "没有现成文档", "没有既定技术约束"]):
        bullet(s, item, 1.10, 2.72 + i * 0.48, 3.8, 15, BLUE)
    rect(s, 6.08, 1.86, 5.75, 2.05, ORANGE_PALE, None, True)
    textbox(s, "原始想法", 6.36, 2.13, 1.8, 0.27, 17, ORANGE, True)
    textbox(s, "开发一个企业工单系统，\n先实现 MVP，\n再迭代开发。", 6.36, 2.62, 4.7, 0.75, 22, BLACK, True)
    rect(s, 6.08, 4.36, 5.75, 1.30, BLUE_PALE, None, True)
    textbox(s, "第一步不是生成代码，\n而是把隐藏选择问出来。", 6.36, 4.72, 5.05, 0.45, 20, BLUE_DARK, True)

    # P06
    s = page(prs, "需求澄清 / 场景", 6, "第一个收敛：这是哪一种工单？")
    choices = [("A", "IT 服务台", "账号、设备、网络"), ("B", "综合服务台", "IT、行政、人事"), ("C", "客户支持/售后", "外部客户问题"), ("D", "运维事件", "告警、故障、变更")]
    for i, (a, h, b) in enumerate(choices):
        x = 0.76 + i * 3.05
        fill = GREEN_PALE if a == "C" else WHITE
        accent = GREEN if a == "C" else BLUE
        rect(s, x, 1.82, 2.64, 1.18, fill, accent if a == "C" else LINE, True, 1.4)
        tag(s, a, x + 0.18, 2.04, 0.38, accent, WHITE if a == "C" else BLUE_PALE)
        textbox(s, h, x + 0.70, 2.02, 1.70, 0.26, 15.5, BLACK, True)
        textbox(s, b, x + 0.24, 2.48, 2.12, 0.20, 11.5, MUTED, False, PP_ALIGN.CENTER)
    textbox(s, "确认结果：C. 客户支持 / 售后", 0.92, 3.38, 5.6, 0.32, 20, GREEN, True)
    table(s, 0.90, 4.10, [2.35, 3.0, 5.7], ["场景", "典型角色", "流程重点"], [
        ["IT 服务台", "内部员工、IT 支持", "内部服务请求与资产问题"],
        ["客户售后", "外部客户、客服、管理员", "客户提交、内部分配、处理闭环"],
        ["运维事件", "值班、研发、运维", "告警响应、影响控制、复盘"],
    ], 0.42, 11.2)

    # P07
    s = page(prs, "需求澄清 / 角色", 7, "角色决定权限，也决定页面")
    card(s, 0.75, 1.86, 3.75, 3.62, "客户", ["开放注册", "创建工单", "查看本人", "公开留言"], BLUE, WHITE, 15.2)
    card(s, 4.78, 1.86, 3.75, 3.62, "客服", ["管理员创建", "查看全部", "按状态筛选", "处理本人负责"], GREEN, WHITE, 15.2)
    card(s, 8.82, 1.86, 3.75, 3.62, "管理员", ["初始化产生", "维护分类/账号", "分配工单", "处理任意工单"], ORANGE, WHITE, 15.2)
    rect(s, 1.04, 5.90, 11.52, 0.54, BLUE_DARK, None, True)
    textbox(s, "客户看本人，客服处理本人负责，管理员管理全部。", 1.14, 6.06, 11.30, 0.23, 18, WHITE, True, PP_ALIGN.CENTER)

    # P08
    s = page(prs, "需求澄清 / MVP 正范围", 8, "MVP 只保留闭环所需能力")
    rect(s, 4.63, 2.20, 3.98, 1.06, BLUE_DARK, None, True)
    textbox(s, "售后工单闭环", 4.88, 2.52, 3.50, 0.28, 24, WHITE, True, PP_ALIGN.CENTER)
    card(s, 0.82, 1.82, 3.25, 3.70, "客户侧", ["注册登录", "创建工单", "查看我的工单", "公开留言"], BLUE, BLUE_PALE, 15)
    card(s, 9.10, 1.82, 3.25, 3.70, "内部侧", ["分类维护", "客服账号", "客户启停", "分配处理关闭"], GREEN, GREEN_PALE, 15)
    timeline(s, ["客户提交", "管理员分配", "客服处理", "工单关闭", "客户查看"], 1.42, 5.92, 10.25, 0.48, GREEN, 11.5)

    # P09
    s = page(prs, "需求澄清 / MVP 边界", 9, "非范围清单，比功能清单更重要")
    rect(s, 0.82, 1.70, 6.68, 4.74, ORANGE_PALE, None, True)
    textbox(s, "MVP 暂不包含", 1.08, 2.00, 2.2, 0.30, 19, ORANGE, True)
    non_scope = ["附件上传", "外部通知", "SLA / 超时提醒", "统计报表", "高级搜索 / 批量 / 导出", "客户确认关闭", "取消 / 退回 / 重开", "密码找回 / MFA", "生产部署 / 数据库迁移"]
    for i, item in enumerate(non_scope):
        bullet(s, item, 1.10 + (i // 5) * 3.00, 2.56 + (i % 5) * 0.45, 2.8, 13.5, ORANGE)
    textbox(s, "不是永远不做，而是不进入 MVP。", 1.08, 5.84, 5.9, 0.24, 16, ORANGE, True)
    table(s, 7.88, 2.05, [2.1, 2.65], ["模糊说法", "明确边界"], [
        ["以后可能要通知", "MVP 不建通知入口"],
        ["以后可能要统计", "迭代 1 再做看板"],
        ["客户最好能确认", "MVP 内部关闭"],
    ], 0.58, 11.0, ORANGE)

    # P10
    s = page(prs, "需求澄清 / 状态", 10, "状态流程越简单，规则越要清楚")
    timeline(s, ["待分配", "处理中", "已解决", "已关闭"], 1.30, 2.02, 10.70, 0.64, BLUE, 16)
    expl = [("等待安排", "管理员分配或介入"), ("内部处理", "客服/管理员回复处理"), ("已有结果", "内部人员认为已解决"), ("结束只读", "MVP 不支持重开")]
    for i, (h, b) in enumerate(expl):
        x = 1.36 + i * 2.70
        textbox(s, h, x, 2.92, 1.95, 0.24, 14, BLACK, True, PP_ALIGN.CENTER)
        textbox(s, b, x - 0.12, 3.26, 2.22, 0.34, 11.2, MUTED, False, PP_ALIGN.CENTER)
    table(s, 1.15, 4.35, [2.2, 3.6, 2.2, 3.6], ["规则", "约定", "规则", "约定"], [
        ["能否跳级", "不允许待分配到已解决", "关闭后留言", "不允许"],
        ["能否回退", "不允许已解决到处理中", "关闭后分配", "不允许"],
        ["分配改状态", "不自动改变", "留言改状态", "不自动改变"],
    ], 0.48, 10.5)

    # P11
    s = page(prs, "需求澄清 / 技术", 11, "技术选型服务于 MVP，也要留下演进出口")
    card(s, 0.82, 1.78, 3.25, 3.85, "当前 MVP", ["原生 HTML/CSS/JS", "FastAPI", "JSON 文件持久化"], BLUE, BLUE_PALE, 15)
    rect(s, 4.72, 2.70, 3.90, 1.20, WHITE, BLUE, True, 1.6)
    textbox(s, "仓储接口层", 5.04, 3.03, 3.28, 0.28, 24, BLUE_DARK, True, PP_ALIGN.CENTER)
    card(s, 9.18, 1.78, 3.25, 3.85, "后续演进", ["真实数据库适配器", "不重写业务流程", "保持接口稳定"], GREEN, GREEN_PALE, 15)
    line(s, 4.08, 3.30, 4.72, 3.30, BLUE, 2.0)
    line(s, 8.62, 3.30, 9.18, 3.30, GREEN, 2.0)
    code_box(s, "页面 -> API/鉴权 -> 业务服务 -> 仓储接口 -> JSON 适配器 / 数据库适配器", 1.18, 5.95, 11.0, 0.50, 13)

    # P12
    s = page(prs, "需求基线 / 上位规范", 12, "第一份正式产物：产品与技术规范 v1.0")
    rect(s, 0.84, 1.82, 3.42, 3.85, BLUE_DARK, None, True)
    textbox(s, "产品与技术\n规范 v1.0", 1.18, 2.72, 2.72, 0.78, 26, WHITE, True, PP_ALIGN.CENTER)
    textbox(s, "已确认", 1.82, 4.34, 1.45, 0.28, 16, GREEN_PALE, True, PP_ALIGN.CENTER)
    questions = ["系统是什么", "MVP 做什么", "MVP 不做什么", "流程规则", "技术方向", "如何验收"]
    for i, q in enumerate(questions):
        x = 4.78 + (i % 2) * 3.55
        y = 1.84 + (i // 2) * 0.90
        mini_card(s, x, y, 3.20, 0.64, q, "", BLUE if i < 5 else GREEN, BLUE_PALE if i < 5 else GREEN_PALE, 15, 10)
    code_box(s, "本规范内容已经用户确认，\n可作为后续 MVP 实施计划的正式依据。", 4.78, 5.00, 6.88, 0.78, 13.5)

    # P13
    s = page(prs, "需求基线 / 落盘", 13, "口头确认不是项目基线")
    card(s, 0.86, 1.78, 5.20, 3.60, "聊天中说过", ["依赖上下文", "后续难以引用", "容易被新讨论覆盖", "不便审查和复用"], ORANGE, ORANGE_PALE, 15)
    card(s, 7.02, 1.78, 5.20, 3.60, "文档中确认", ["路径明确", "版本明确", "状态明确", "可被后续文档引用"], GREEN, GREEN_PALE, 15)
    code_box(s, "docs/00-项目基线与记录/\n  企业售后工单系统-产品与技术规范-v1.0.md", 1.52, 5.86, 10.18, 0.62, 13)

    # P14
    s = page(prs, "需求基线 / PRD", 14, "PRD 讲清楚产品目标和用户价值")
    items = ["产品背景", "产品目标", "目标用户", "核心场景", "MVP 非目标", "产品验收目标"]
    for i, item in enumerate(items):
        x = 0.90 + (i % 3) * 3.92
        y = 1.82 + (i // 3) * 0.80
        mini_card(s, x, y, 3.38, 0.58, item, "", BLUE, BLUE_PALE, 15, 10)
    table(s, 1.04, 3.88, [1.40, 9.80], ["编号", "MVP 目标示例"], [
        ["G-01", "为个人客户提供统一的在线售后问题提交与进展查看入口。"],
        ["G-03", "使负责客服能够基于工单与客户沟通，并推进问题至关闭。"],
        ["G-05", "验证可演进的数据访问架构，为后续真实数据库接入提供基础。"],
    ], 0.52, 12)

    # P15
    s = page(prs, "需求基线 / SRS", 15, "SRS 把“能做”变成“必须如何做”")
    rect(s, 0.86, 1.92, 3.54, 1.08, BLUE_DARK, None, True)
    textbox(s, "PRD：\n客户应能创建工单。", 1.12, 2.15, 3.05, 0.48, 18, WHITE, True, PP_ALIGN.CENTER)
    qs = [("谁能创建？", "已登录且启用客户"), ("创建什么？", "标题、描述、分类"), ("字段规则？", "标题 1-100 字"), ("分类校验？", "必须启用"), ("创建后状态？", "待分配"), ("能否修改？", "MVP 不允许")]
    for i, (h, b) in enumerate(qs):
        x = 4.78 + (i % 2) * 3.42
        y = 1.58 + (i // 2) * 0.88
        mini_card(s, x, y, 3.16, 0.72, h, b, GREEN, GREEN_PALE, 12.5, 9.8)
    table(s, 0.92, 4.72, [1.72, 9.45], ["编号", "需求片段"], [
        ["SRS-TKT-001", "仅已登录且启用的客户能够创建工单。"],
        ["SRS-TKT-003", "保存前必须再次验证所选分类仍为有效状态。"],
        ["SRS-TKT-004", "创建成功时，状态为待分配，负责人为空。"],
    ], 0.42, 10.6, GREEN)

    # P16
    s = page(prs, "需求基线 / 页面交互", 16, "用户不是操作需求条目，而是操作页面")
    card(s, 0.82, 1.72, 3.45, 2.42, "客户门户", ["登录页", "客户注册页", "我的工单页", "新建工单页", "客户详情页"], BLUE, BLUE_PALE, 13.5)
    card(s, 4.60, 1.72, 3.45, 2.42, "内部后台", ["内部列表页", "内部详情页", "分类管理页", "客服账号页", "客户账号页"], GREEN, GREEN_PALE, 13.5)
    table(s, 0.95, 4.55, [2.2, 4.0, 4.5], ["页面区域", "展示与操作", "规则"], [
        ["问题分类", "下拉选择有效分类", "无有效分类时禁止提交"],
        ["标题", "必填输入框", "1-100 字"],
        ["问题描述", "多行文本", "1-4000 字"],
        ["提交按钮", "创建工单", "成功后跳转详情页"],
    ], 0.40, 10.2)
    rect(s, 8.60, 1.72, 3.48, 2.42, WHITE, LINE, True)
    textbox(s, "创建成功流转", 8.82, 2.00, 2.8, 0.26, 16, BLACK, True)
    textbox(s, "填写提交\n-> 后端校验\n-> 创建待分配工单\n-> 跳转详情\n-> 提示创建成功", 8.88, 2.48, 2.9, 1.10, 13, TEXT)

    # P17
    s = page(prs, "需求基线 / 追踪关系", 17, "重要规则不能只写一次")
    rect(s, 3.68, 1.70, 5.96, 0.75, BLUE_DARK, None, True)
    textbox(s, "客户不展示具体客服账号或负责人名称", 3.82, 1.94, 5.68, 0.25, 19, WHITE, True, PP_ALIGN.CENTER)
    nodes = [("SRS", "数据边界"), ("页面", "详情不显示"), ("API", "响应不返回"), ("测试", "TC-SEC-004"), ("实施", "M5 安全复核")]
    for i, (h, b) in enumerate(nodes):
        x = 0.82 + i * 2.43
        rect(s, x, 3.38, 1.90, 1.02, GREEN_PALE if i == 3 else BLUE_PALE, None, True)
        textbox(s, h, x + 0.10, 3.62, 1.70, 0.20, 15, GREEN if i == 3 else BLUE, True, PP_ALIGN.CENTER)
        textbox(s, b, x + 0.10, 3.94, 1.70, 0.20, 11.5, TEXT, False, PP_ALIGN.CENTER)
        if i < len(nodes) - 1:
            textbox(s, ">", x + 1.90, 3.70, 0.52, 0.22, 15, MUTED, True, PP_ALIGN.CENTER)
    textbox(s, "客户需要看到处理进展，但不需要看到内部负责人账号。", 1.36, 5.52, 10.48, 0.28, 17, BLUE_DARK, True, PP_ALIGN.CENTER)

    # P18
    s = page(prs, "需求基线 / 文档体系", 18, "不同文档，回答不同层面的问题")
    table(s, 0.72, 1.66, [2.55, 5.42, 3.55], ["文档", "回答的问题", "主要使用者"], [
        ["产品与技术规范", "项目总体边界是什么？", "全体"],
        ["PRD", "为什么做，为谁做？", "产品/业务"],
        ["SRS", "系统必须表现出什么行为？", "开发/测试"],
        ["页面说明", "用户在哪些页面完成操作？", "前端/产品"],
        ["技术/数据/API", "如何实现、保存和通信？", "研发/测试"],
        ["测试方案", "怎样证明符合基线？", "测试/验收"],
        ["实施计划", "按什么顺序开发交付？", "项目/开发"],
    ], 0.48, 10.5)

    # P19
    s = page(prs, "设计基线 / 技术设计", 19, "架构不是图好看，而是规则有地方落")
    layers = ["静态前端", "FastAPI 接口层", "应用服务层", "领域规则层", "仓储协议层", "JSON 适配器"]
    for i, layer in enumerate(layers):
        rect(s, 0.92, 1.72 + i * 0.62, 4.60, 0.42, BLUE_PALE if i < 5 else GREEN_PALE, LINE, True)
        textbox(s, layer, 1.12, 1.84 + i * 0.62, 4.15, 0.16, 13.5, BLUE_DARK if i < 5 else GREEN, True, PP_ALIGN.CENTER)
    table(s, 6.10, 1.78, [3.0, 3.2], ["规则", "落点"], [
        ["客户本人权限", "领域规则 / 服务层"],
        ["不返回密码哈希", "响应模型 / 接口层"],
        ["分配和记录一次保存", "服务 / 仓储"],
        ["JSON 原子替换", "存储基础设施"],
    ], 0.55, 11.0)

    # P20
    s = page(prs, "设计基线 / 权限", 20, "按钮隐藏不是权限控制")
    timeline(s, ["Cookie", "会话", "用户状态", "角色", "资源归属", "工单状态", "业务"], 0.76, 1.78, 11.78, 0.50, BLUE, 10.5)
    table(s, 0.82, 3.10, [2.3, 2.0, 2.0, 2.0, 2.7], ["操作", "客户", "客服", "管理员", "规则"], [
        ["创建工单", "允许", "不允许", "不允许", "仅客户"],
        ["查看全部", "不允许", "允许", "允许", "内部可见"],
        ["发送留言", "本人工单", "本人负责", "任意工单", "未关闭"],
        ["分配工单", "不允许", "不允许", "允许", "管理员"],
        ["推进状态", "不允许", "本人负责", "任意工单", "线性推进"],
    ], 0.48, 10.4, BLUE)

    # P21
    s = page(prs, "设计基线 / 数据模型", 21, "数据模型决定可追溯能力")
    entities = ["User", "Category", "Ticket", "Message", "AuditLog", "Session"]
    for i, e in enumerate(entities):
        x = 0.98 + (i % 3) * 2.00
        y = 1.74 + (i // 3) * 0.72
        mini_card(s, x, y, 1.62, 0.52, e, "", BLUE if e != "AuditLog" else GREEN, BLUE_PALE if e != "AuditLog" else GREEN_PALE, 13, 10)
    table(s, 0.88, 3.55, [2.9, 6.5], ["Ticket 字段", "含义"], [
        ["id", "UUID v4 工单标识"],
        ["title / description", "标题 1-100 字，描述 1-4000 字"],
        ["category_name_snapshot", "创建时分类名称快照"],
        ["customer_user_id", "创建客户"],
        ["status / assignee_user_id", "当前状态与负责客服"],
    ], 0.42, 10.5)
    rect(s, 8.42, 1.74, 3.35, 1.36, ORANGE_PALE, None, True)
    textbox(s, "为什么需要快照？\n分类后续可以改名或停用，\n历史工单仍显示提交时分类名称。", 8.66, 2.00, 2.92, 0.72, 13, ORANGE)

    # P22
    s = page(prs, "设计基线 / JSON 存储", 22, "JSON 是当前实现，仓储接口是未来出口")
    code_box(s, '{\n  "schema_version": 1,\n  "meta": {...},\n  "users": [],\n  "categories": [],\n  "tickets": [],\n  "messages": [],\n  "audit_logs": [],\n  "sessions": []\n}', 0.88, 1.74, 4.72, 3.46, 11)
    card(s, 6.20, 1.74, 2.70, 3.46, "写入保护", ["写入锁", "读取快照", "内存变更", "临时文件", "原子替换"], BLUE, BLUE_PALE, 13.5)
    card(s, 9.34, 1.74, 2.70, 3.46, "演进边界", ["服务层依赖协议", "JSON 是适配器", "未来数据库替换"], GREEN, GREEN_PALE, 13.5)
    textbox(s, "服务层只依赖仓储协议，不直接读写 JSON 文件。", 1.10, 5.88, 11.05, 0.28, 17, BLUE_DARK, True, PP_ALIGN.CENTER)

    # P23
    s = page(prs, "设计基线 / API", 23, "API 是页面和业务服务之间的合同")
    groups = ["/auth", "/categories", "/customer", "/internal", "/admin"]
    for i, g in enumerate(groups):
        tag(s, g, 0.92 + i * 2.32, 1.72, 1.55, BLUE if i < 3 else GREEN, BLUE_PALE if i < 3 else GREEN_PALE, 12)
    code_box(s, 'POST /api/v1/customer/tickets\n\n{\n  "title": "无法正常使用服务",\n  "description": "登录后页面无法继续操作。",\n  "category_id": "550e8400-..."\n}', 0.90, 2.58, 5.50, 2.90, 10.5)
    code_box(s, '{\n  "error": {\n    "code": "VALIDATION_ERROR",\n    "message": "输入内容不符合要求。",\n    "field_errors": {\n      "title": "标题长度必须为 1 至 100 个字符。"\n    }\n  }\n}', 6.76, 2.58, 5.50, 2.90, 10.0)

    # P24
    s = page(prs, "设计基线 / 三件套", 24, "三份设计文档共同构成实现蓝图")
    card(s, 0.92, 2.02, 3.28, 2.60, "技术设计", ["结构与机制", "认证与授权", "存储策略"], BLUE, BLUE_PALE)
    card(s, 5.02, 2.02, 3.28, 2.60, "数据模型", ["事实与约束", "实体与关联", "历史快照"], GREEN, GREEN_PALE)
    card(s, 9.12, 2.02, 3.28, 2.60, "API 设计", ["交互与边界", "权限与错误", "字段暴露"], ORANGE, ORANGE_PALE)
    textbox(s, "MVP 设计基线", 3.80, 5.32, 5.72, 0.38, 24, BLUE_DARK, True, PP_ALIGN.CENTER)
    textbox(s, "系统怎么组织  |  事实怎么保存  |  前后端怎么约定", 2.10, 6.02, 9.10, 0.24, 14, MUTED, False, PP_ALIGN.CENTER)

    # P25
    s = page(prs, "验证前置 / 一致性审查", 25, "开发前先问：需求都落到设计了吗？")
    textbox(s, "审查对象：PRD / SRS / 页面说明 / 技术设计 / 数据模型 / API 设计", 0.92, 1.62, 11.5, 0.26, 15, BLUE_DARK, True)
    tags = ["88 项 SRS 需求", "角色权限", "状态流程", "敏感字段隔离", "JSON 持久化", "MVP 范围"]
    for i, t in enumerate(tags):
        tag(s, t, 0.96 + (i % 3) * 3.78, 2.18 + (i // 3) * 0.50, 2.60, GREEN if i == 0 else BLUE, GREEN_PALE if i == 0 else BLUE_PALE, 10.5)
    table(s, 0.86, 3.55, [2.2, 2.1, 2.8, 3.4, 1.0], ["SRS 范围", "主题", "页面落实", "技术/数据/API", "结论"], [
        ["TKT-006~011", "工单查看", "列表与详情", "权限校验、字段隔离", "覆盖"],
        ["SEC-001~004", "授权", "角色入口", "服务端授权、资源级检查", "覆盖"],
        ["DAT-001~007", "持久化", "无直接界面", "仓储接口、store.json", "覆盖"],
    ], 0.46, 9.2, BLUE)

    # P26
    s = page(prs, "验证前置 / 发现项", 26, "发现问题不是坏事，未关闭才危险")
    card(s, 0.82, 1.78, 3.65, 2.28, "F-01", ["客户详情负责人展示歧义", "影响字段隔离"], ORANGE, ORANGE_PALE, 12.5)
    card(s, 4.84, 1.78, 3.65, 2.28, "F-02", ["建单成功跳转不确定", "影响页面行为"], ORANGE, ORANGE_PALE, 12.5)
    card(s, 8.86, 1.78, 3.65, 2.28, "F-03", ["技术方案建议语气", "影响设计约束"], ORANGE, ORANGE_PALE, 12.5)
    table(s, 0.96, 4.74, [2.0, 9.4], ["发现", "修订结果"], [
        ["F-01", "客户详情不展示具体客服账号或负责人名称"],
        ["F-02", "创建成功后跳转新建工单详情页"],
        ["F-03", "统一改为确定性设计约束"],
    ], 0.43, 10.8, GREEN)

    # P27
    s = page(prs, "验证前置 / 测试方案", 27, "测试不是最后补充，而是提前定义“怎样算做对”")
    card(s, 0.84, 1.70, 3.48, 3.12, "测试目标", ["业务闭环", "权限边界", "状态数据一致", "凭证安全", "本地持久化", "页面反馈"], BLUE, BLUE_PALE, 12.8)
    card(s, 4.70, 1.70, 3.48, 3.12, "测试层次", ["规则单元", "仓储服务", "API 集成", "页面联调", "业务验收", "文档/代码审查"], GREEN, GREEN_PALE, 12.8)
    table(s, 8.38, 1.72, [1.20, 1.50, 1.35], ["用例", "场景", "预期"], [
        ["AUTH-004", "三类登录", "进入首页"],
        ["TKT-005", "他人工单", "拒绝访问"],
        ["STS-004", "跳级回退", "拒绝变更"],
        ["DAT-002", "重启读取", "数据存在"],
    ], 0.55, 7.8, BLUE)

    # P28
    s = page(prs, "验证前置 / E2E", 28, "验收要覆盖主流程，也要覆盖风险点")
    e2e = [
        ("E2E-01", "初始化与建单", "客户建单待分配且可查看"),
        ("E2E-02", "分配处理关闭", "客户可看进展，关闭后只读"),
        ("E2E-03", "越权验证", "非负责人不可写，其他客户不可读"),
        ("E2E-04", "禁用客户", "会话失效，历史保留"),
        ("E2E-05", "重启读取", "账号工单留言记录完整"),
    ]
    for i, (code, title, body) in enumerate(e2e):
        x = 0.82 + (i % 3) * 4.02
        y = 1.82 + (i // 3) * 1.62
        mini_card(s, x, y, 3.45, 1.20, f"{code} {title}", body, GREEN if i in [0, 1] else BLUE, GREEN_PALE if i in [0, 1] else BLUE_PALE, 14, 11.5)

    # P29
    s = page(prs, "验证前置 / 质量门禁", 29, "有些缺陷不能靠“后续优化”处理")
    redlines = ["越权访问", "敏感数据泄露", "明文密码", "状态违规写入", "数据损坏", "闭环不可完成"]
    for i, item in enumerate(redlines):
        tag(s, item, 0.98 + (i % 3) * 3.72, 1.76 + (i // 3) * 0.58, 2.25, ORANGE, ORANGE_PALE, 11.5)
    table(s, 1.06, 3.48, [2.0, 5.2, 3.8], ["等级", "示例", "验收处理"], [
        ["严重", "越权读取、明文密码、JSON 损坏", "必须修复并复测"],
        ["高", "状态可跳级、退出后仍可访问", "必须修复并复测"],
        ["中", "重要交互反馈缺失", "原则上修复"],
        ["低", "不影响流程的展示问题", "记录后验收"],
    ], 0.50, 10.5, ORANGE)

    # P30
    s = page(prs, "实施计划 / 里程碑", 30, "计划不是排期表，而是风险处理顺序")
    milestones = ["M0\n工程骨架", "M1\n数据/安全/会话", "M2\n管理配置", "M3\n客户工单", "M4\n内部闭环", "M5\n安全联调", "M6\n验收交付"]
    timeline(s, milestones, 0.78, 2.28, 11.80, 0.70, BLUE, 11)
    details = ["应用启动", "基础能力", "分类账号", "建单详情", "分配处理", "字段安全", "全量验收"]
    for i, d in enumerate(details):
        textbox(s, d, 0.78 + i * 1.70, 3.24, 1.52, 0.22, 11, MUTED, False, PP_ALIGN.CENTER)
    textbox(s, "基础能力先固化，业务闭环再纵向集成，最后做安全与验收复核。", 1.20, 5.42, 10.90, 0.28, 18, BLUE_DARK, True, PP_ALIGN.CENTER)

    # P31
    s = page(prs, "实施计划 / 拆解逻辑", 31, "页面不是最小风险单元")
    card(s, 0.88, 1.70, 5.30, 3.90, "直接按页面开写", ["登录页临时会话", "建单页临时校验", "详情页临时权限", "管理页临时审计", "最后补权限和测试"], ORANGE, ORANGE_PALE, 14)
    card(s, 6.88, 1.70, 5.30, 3.90, "推荐路径", ["先固化数据/仓储/安全/会话", "再完成管理配置", "再完成客户建单", "再完成内部处理", "最后字段隔离和验收"], BLUE, BLUE_PALE, 14)

    # P32
    s = page(prs, "实施计划 / 测试驱动", 32, "高风险规则，适合先写测试")
    card(s, 0.82, 1.65, 3.48, 3.00, "适合测试先行", ["状态转换", "权限判断", "资源归属", "会话失效", "字段隔离", "JSON 原子写入"], GREEN, GREEN_PALE, 12.5)
    table(s, 4.55, 1.78, [1.38, 1.38, 1.10], ["当前", "下一状态", "预期"], [
        ["待分配", "处理中", "允许"],
        ["处理中", "已解决", "允许"],
        ["待分配", "已解决", "拒绝"],
        ["已解决", "处理中", "拒绝"],
        ["已关闭", "任意", "拒绝"],
    ], 0.42, 8.6, GREEN)
    code_box(s, "客户工单详情响应不得包含：\n- assignee_user_id\n- assignee.username\n- audit_logs\n- password_hash\n- token_hash", 8.72, 1.80, 3.46, 2.48, 10.4)
    rect(s, 0.82, 4.82, 11.34, 0.70, WHITE, LINE, True)
    rule(s, 1.00, 5.00, 0.46, BLUE, 0.045)
    textbox(s, "不必强制严格 TDD", 1.62, 4.98, 2.58, 0.22, 14, BLACK, True)
    textbox(s, "普通页面排版、CSS 视觉细节、一般提示文案微调", 4.42, 4.99, 7.35, 0.20, 12.2, TEXT)

    # P33
    s = page(prs, "实施计划 / 启动检查", 33, "不是文档写完了，而是条件满足了")
    checks = [
        ("已完成", "产品与技术规范已经确认", GREEN),
        ("已完成", "PRD、SRS 和页面交互说明已经确认", GREEN),
        ("已完成", "技术设计、数据模型和 API 设计已经确认", GREEN),
        ("已完成", "需求与设计一致性审查通过", GREEN),
        ("已完成", "测试与验收方案已经确认", GREEN),
        ("待确认", "MVP 实施计划最终确认", BLUE),
        ("待启动", "代码工程初始化", ORANGE),
    ]
    for i, (status, txt, accent) in enumerate(checks):
        y = 1.56 + i * 0.58
        tag(s, status, 1.00, y, 1.00, accent, GREEN_PALE if accent == GREEN else BLUE_PALE if accent == BLUE else ORANGE_PALE, 10)
        textbox(s, txt, 2.32, y + 0.04, 7.6, 0.22, 15, TEXT, True)
    rect(s, 1.08, 6.04, 10.94, 0.50, BLUE_DARK, None, True)
    textbox(s, "当前项目尚未开始代码开发，正在形成开发前实施计划。", 1.18, 6.19, 10.72, 0.19, 16, WHITE, True, PP_ALIGN.CENTER)

    # P34
    s = page(prs, "Agent 协作 / Agent 作用", 34, "Agent 的价值，是加快分析和产物衔接")
    rect(s, 4.85, 2.70, 3.55, 1.08, BLUE_DARK, None, True)
    textbox(s, "Agent 协作", 5.14, 3.02, 2.96, 0.28, 24, WHITE, True, PP_ALIGN.CENTER)
    left = ["用户输入", "业务选择", "确认批准"]
    right = ["澄清问题", "规范文档", "设计文档", "追踪矩阵", "测试方案", "实施计划"]
    card(s, 0.86, 1.95, 2.80, 2.55, "人提供", left, GREEN, GREEN_PALE, 14)
    card(s, 9.42, 1.58, 2.90, 3.40, "产物输出", right, BLUE, BLUE_PALE, 12.5)
    line(s, 3.66, 3.24, 4.85, 3.24, GREEN, 2.0)
    line(s, 8.40, 3.24, 9.42, 3.24, BLUE, 2.0)
    textbox(s, "分析、整理、生成、检查、拆解", 3.28, 5.55, 6.92, 0.27, 18, BLUE_DARK, True, PP_ALIGN.CENTER)

    # P35
    s = page(prs, "Agent 协作 / 人的责任", 35, "Agent 可以推进工作，人必须批准方向")
    card(s, 0.84, 1.70, 5.45, 4.10, "人负责", ["选择业务场景", "确认客户与账号模式", "确认处理组织", "确认范围取舍", "确认技术约束", "批准基线", "判断是否进入下一阶段"], GREEN, GREEN_PALE, 12.8)
    card(s, 7.00, 1.70, 5.00, 4.10, "Agent 协助", ["提供选项", "整理材料", "生成文档", "检查一致性", "拆解计划"], BLUE, BLUE_PALE, 13.5)

    # P36
    s = page(prs, "Agent 协作 / 过程纪律", 36, "Agent 协作也需要过程纪律")
    lessons = ["确认后要落盘", "建议措辞要定稿", "发现项要关闭", "非范围不要偷做", "测试标准不能后补"]
    for i, item in enumerate(lessons):
        tag(s, item, 0.88 + (i % 3) * 4.02, 1.72 + (i // 3) * 0.58, 2.62, ORANGE, ORANGE_PALE, 10.8)
    table(s, 1.05, 3.56, [4.5, 6.5], ["现象", "修正"], [
        ["规范确认后未立即保存", "补充落盘为正式基线文档"],
        ["客户负责人展示存在歧义", "修订为客户详情不展示负责人"],
        ["技术设计仍有建议措辞", "改为确定性设计约束"],
    ], 0.50, 10.8, ORANGE)

    # P37
    s = page(prs, "Agent 协作 / 可复用模式", 37, "用同一套节奏处理下一个项目")
    steps = ["问清场景", "确认 MVP", "写清非范围", "固化基线", "拆需求交互", "展开设计", "先定验收", "再排计划"]
    for i, step in enumerate(steps):
        x = 0.78 + i * 1.52
        y = 4.95 - i * 0.33
        rect(s, x, y, 1.20, 0.58, BLUE_PALE if i < 6 else GREEN_PALE, None, True)
        textbox(s, step, x + 0.05, y + 0.17, 1.10, 0.16, 10.2, BLUE if i < 6 else GREEN, True, PP_ALIGN.CENTER)
    textbox(s, "每一步都有确认，每一步都有产物。", 1.36, 5.86, 10.48, 0.30, 20, BLUE_DARK, True, PP_ALIGN.CENTER)

    # P38
    s = page(prs, "收束 / 总结", 38, "规格清楚，Agent 才能跑得稳", key=True)
    support = ["范围有边界", "实现有依据", "测试有证据", "变更可控制", "验收可判断"]
    for i, item in enumerate(support):
        tag(s, item, 1.02 + i * 2.36, 2.18, 1.60, GREEN if i == 2 else BLUE, GREEN_PALE if i == 2 else BLUE_PALE, 11)
    table(s, 1.02, 3.34, [5.35, 5.35], ["没有基线的快", "有基线的快"], [
        ["快速生成，但假设不清", "快速执行，边界清楚"],
        ["页面能跑，但权限不稳", "流程可用，规则可测"],
        ["需求变了难追踪", "变更影响可定位"],
        ["验收时再争论", "开发前已定义通过标准"],
    ], 0.48, 11.3, BLUE_DARK)

    # P39
    s = page(prs, "收束 / 课堂练习", 39, "练习：把一句想法变成 MVP 边界")
    card(s, 0.82, 1.66, 4.35, 4.85, "练习步骤", ["写一句原始想法", "设计 5 个澄清问题", "写 3 项 MVP 必做", "写 5 项暂不做", "定义 1 条 E2E 验收路径"], BLUE, BLUE_PALE, 14)
    table(s, 5.58, 1.70, [2.55, 4.12], ["项目", "填写"], [
        ["原始想法", ""],
        ["目标用户", ""],
        ["5 个澄清问题", ""],
        ["MVP 必做 3 项", ""],
        ["MVP 暂不做 5 项", ""],
        ["E2E 验收路径", ""],
    ], 0.55, 10.8, GREEN)
    textbox(s, "示例：内部知识库系统 -> 创建知识条目、搜索条目、管理员审核", 0.95, 6.68, 11.6, 0.20, 12, MUTED)

    # P40
    s = page(prs, "收束 / 讨论", 40, "讨论：你的团队最容易跳过哪一步？")
    questions = [("需求澄清", "一开始默认理解一致"), ("非范围确认", "只写做什么，不写暂不做"), ("设计追踪", "需求、页面、数据、API 缺少映射"), ("测试前置", "代码后才讨论完成标准"), ("变更确认", "开发中悄悄扩大范围")]
    for i, (h, b) in enumerate(questions):
        x = 0.82 + (i % 3) * 4.02
        y = 1.78 + (i // 3) * 1.16
        mini_card(s, x, y, 3.42, 0.86, h, b, ORANGE if i in [1, 4] else BLUE, ORANGE_PALE if i in [1, 4] else BLUE_PALE, 14, 10.8)
    rect(s, 1.02, 5.64, 11.25, 0.72, BLUE_DARK, None, True)
    textbox(s, "Agent 加快执行，人负责判断；基线越清楚，协作越可靠。", 1.16, 5.88, 10.98, 0.23, 19, WHITE, True, PP_ALIGN.CENTER)
    textbox(s, "Q&A", 11.18, 6.68, 1.0, 0.24, 15, BLUE, True, PP_ALIGN.RIGHT)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(str(OUT))
    print(f"slides={len(prs.slides)}")
    return OUT


if __name__ == "__main__":
    build()
