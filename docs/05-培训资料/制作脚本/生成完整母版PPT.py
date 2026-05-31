# -*- coding: utf-8 -*-
"""Generate the full training-deck master presentation."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parent.parent / "03-传统软件开发模式与AI-Agent时代的软件开发模式-完整母版.pptx"

FONT = "Microsoft YaHei"
BLUE = RGBColor(0x14, 0x64, 0xA5)
BLUE_DARK = RGBColor(0x0C, 0x4F, 0x87)
BLUE_LIGHT = RGBColor(0xEE, 0xF5, 0xFB)
GREEN = RGBColor(0x20, 0xA3, 0x6A)
GREEN_LIGHT = RGBColor(0xE9, 0xF7, 0xF0)
ORANGE = RGBColor(0xF2, 0x8C, 0x28)
ORANGE_LIGHT = RGBColor(0xFE, 0xF3, 0xE8)
BLACK = RGBColor(0x11, 0x11, 0x11)
TEXT = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0xD9, 0xE2, 0xEC)
GRAY_TEXT = RGBColor(0x68, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_PT = 28
BODY_PT = 18
META_PT = 10
W = 13.333
H = 7.5


def I(value):
    return Inches(value)


def set_east_asian(run):
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    if ea is None:
        ea = OxmlElement("a:ea")
        rpr.append(ea)
    ea.set("typeface", FONT)


def style_run(run, size=BODY_PT, color=TEXT, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    set_east_asian(run)


def text_box(slide, text, x, y, w, h, size=BODY_PT, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE, margin=0.03):
    shape = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = I(margin)
    tf.margin_right = I(margin)
    tf.margin_top = I(margin)
    tf.margin_bottom = I(margin)
    lines = text.split("\n")
    for index, line_text in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = line_text
        p.alignment = align
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        p.line_spacing = 1.08
        for run in p.runs:
            style_run(run, size, color, bold)
    return shape


def box(slide, x, y, w, h, fill=WHITE, line=GRAY, rounded=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, I(x), I(y), I(w), I(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def fill_bar(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x), I(y), I(w), I(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def connector(slide, x1, y1, x2, y2, color=BLUE, width=2):
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    return shape


def circle(slide, x, y, d, fill=WHITE, line=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, I(x), I(y), I(d), I(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.5)
    return shape


def chevron(slide, x, y, w=0.22, h=0.22, color=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, I(x), I(y), I(w), I(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def title(slide, value, section, number):
    text_box(slide, value, 0.62, 0.32, 10.8, 0.52, TITLE_PT, BLACK, True)
    fill_bar(slide, 0.63, 0.93, 1.0, 0.045, BLUE)
    tag = box(slide, 11.38, 0.38, 1.25, 0.32, BLUE_LIGHT, BLUE_LIGHT)
    text_box(slide, section, 11.42, 0.39, 1.17, 0.28, 10, BLUE_DARK, True, PP_ALIGN.CENTER)
    fill_bar(slide, 0.63, 7.18, 0.32, 0.035, BLUE)
    text_box(slide, "传统软件开发模式与 AI Agent 协作开发", 1.02, 7.04, 4.6, 0.24,
             META_PT, GRAY_TEXT)
    text_box(slide, f"{number:02d} / 22", 11.72, 7.04, 0.9, 0.24,
             META_PT, GRAY_TEXT, False, PP_ALIGN.RIGHT)


def slide_base(prs, page_title, section, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    title(slide, page_title, section, number)
    return slide


def mini_icon_container(slide, x, y, color=BLUE):
    circle(slide, x, y, 0.48, BLUE_LIGHT, BLUE_LIGHT)
    return color


def icon_building(slide, x, y, color=BLUE):
    mini_icon_container(slide, x, y, color)
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x + 0.14), I(y + 0.13), I(0.20), I(0.25))
    b.fill.background()
    b.line.color.rgb = color
    b.line.width = Pt(1.2)
    for dx in (0.18, 0.27):
        for dy in (0.18, 0.26):
            win = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x + dx), I(y + dy), I(0.035), I(0.035))
            win.fill.solid()
            win.fill.fore_color.rgb = color
            win.line.fill.background()


def icon_document(slide, x, y, color=BLUE):
    mini_icon_container(slide, x, y, color)
    d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x + 0.15), I(y + 0.10), I(0.19), I(0.28))
    d.fill.solid()
    d.fill.fore_color.rgb = WHITE
    d.line.color.rgb = color
    d.line.width = Pt(1.1)
    for dy in (0.18, 0.24, 0.30):
        connector(slide, x + 0.18, y + dy, x + 0.30, y + dy, color, 1)


def icon_network(slide, x, y, color=BLUE):
    mini_icon_container(slide, x, y, color)
    nodes = [(x + 0.23, y + 0.14), (x + 0.14, y + 0.29), (x + 0.32, y + 0.29)]
    connector(slide, nodes[0][0], nodes[0][1], nodes[1][0], nodes[1][1], color, 1)
    connector(slide, nodes[0][0], nodes[0][1], nodes[2][0], nodes[2][1], color, 1)
    connector(slide, nodes[1][0], nodes[1][1], nodes[2][0], nodes[2][1], color, 1)
    for nx, ny in nodes:
        circle(slide, nx - 0.035, ny - 0.035, 0.07, color, color)


def icon_check(slide, x, y, color=GREEN):
    mini_icon_container(slide, x, y, color)
    circle(slide, x + 0.12, y + 0.12, 0.24, WHITE, color)
    connector(slide, x + 0.18, y + 0.24, x + 0.22, y + 0.29, color, 1.6)
    connector(slide, x + 0.22, y + 0.29, x + 0.31, y + 0.18, color, 1.6)


def icon_shield(slide, x, y, color=ORANGE):
    mini_icon_container(slide, x, y, color)
    shield = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, I(x + 0.13), I(y + 0.12), I(0.23), I(0.25))
    shield.fill.background()
    shield.line.color.rgb = color
    shield.line.width = Pt(1.3)
    text_box(slide, "!", x + 0.19, y + 0.15, 0.11, 0.16, 12, color, True, PP_ALIGN.CENTER)


def icon_agent(slide, x, y, color=BLUE):
    mini_icon_container(slide, x, y, color)
    circle(slide, x + 0.18, y + 0.18, 0.12, color, color)
    for nx, ny in ((x + 0.11, y + 0.10), (x + 0.34, y + 0.13), (x + 0.13, y + 0.34), (x + 0.34, y + 0.33)):
        connector(slide, x + 0.24, y + 0.24, nx, ny, color, 1)
        circle(slide, nx - 0.025, ny - 0.025, 0.05, color, color)


def icon_target(slide, x, y, color=BLUE):
    mini_icon_container(slide, x, y, color)
    circle(slide, x + 0.12, y + 0.12, 0.24, WHITE, color)
    circle(slide, x + 0.18, y + 0.18, 0.12, WHITE, color)
    circle(slide, x + 0.215, y + 0.215, 0.05, color, color)


def pill(slide, text, x, y, w, color=BLUE, fill=BLUE_LIGHT):
    box(slide, x, y, w, 0.4, fill, fill)
    text_box(slide, text, x + 0.06, y + 0.04, w - 0.12, 0.30, 14, color, True, PP_ALIGN.CENTER)


def header_card(slide, x, y, w, h, header, lines, accent=BLUE, fill=WHITE, icon=None):
    box(slide, x, y, w, h, fill, GRAY)
    fill_bar(slide, x, y, 0.06, h, accent)
    if icon:
        icon(slide, x + 0.17, y + 0.17, accent)
        head_x = x + 0.76
        head_w = w - 0.94
    else:
        head_x = x + 0.20
        head_w = w - 0.38
    text_box(slide, header, head_x, y + 0.17, head_w, 0.36, BODY_PT, BLACK, True)
    for index, line in enumerate(lines):
        yy = y + 0.72 + index * 0.43
        circle(slide, x + 0.22, yy + 0.10, 0.07, accent, accent)
        text_box(slide, line, x + 0.39, yy, w - 0.62, 0.30, BODY_PT, TEXT)


def small_stat(slide, x, y, w, label, value, accent):
    box(slide, x, y, w, 0.98, WHITE, GRAY)
    fill_bar(slide, x, y, 0.06, 0.98, accent)
    text_box(slide, label, x + 0.18, y + 0.16, w - 0.3, 0.28, 16, GRAY_TEXT, False)
    text_box(slide, value, x + 0.18, y + 0.49, w - 0.3, 0.32, BODY_PT, BLACK, True)


def add_notes(slide, text):
    # python-pptx does not offer a stable public API for notes; the source script retains notes.
    return text


def build_deck():
    prs = Presentation()
    prs.slide_width = I(W)
    prs.slide_height = I(H)
    prs.core_properties.title = "传统软件开发模式与 AI Agent 时代的软件开发模式"
    prs.core_properties.subject = "完整培训母版"
    prs.core_properties.author = "Codex"

    # P01
    s = slide_base(prs, "传统软件开发模式与 AI Agent 时代的软件开发模式", "导入", 1)
    text_box(s, "软件协作方式正在变化", 0.68, 1.68, 5.4, 0.47, BODY_PT, BLUE_DARK, True)
    text_box(s, "面向产品、项目与研发管理人员", 0.68, 2.28, 5.5, 0.42, BODY_PT, TEXT)
    pill(s, "单位名称", 0.68, 4.60, 1.55)
    pill(s, "讲师姓名", 2.42, 4.60, 1.55)
    pill(s, "日期", 4.17, 4.60, 1.25)
    steps = [("阶段管理", icon_building, BLUE), ("持续反馈", icon_network, BLUE), ("Agent 协作", icon_agent, GREEN)]
    for i, (label, icon, color) in enumerate(steps):
        x = 7.02 + i * 1.83
        icon(s, x + 0.40, 2.13, color)
        box(s, x, 2.85, 1.36, 0.86, BLUE_LIGHT if color == BLUE else GREEN_LIGHT, color)
        text_box(s, label, x + 0.08, 3.10, 1.20, 0.32, BODY_PT, color, True, PP_ALIGN.CENTER)
        if i < 2:
            chevron(s, x + 1.51, 3.15, 0.20, 0.20, BLUE)
    icon_building(s, 11.70, 5.53, BLUE)
    add_notes(s, "本课程讨论软件开发模式如何变化，而非某一个 AI 工具的操作技巧。")

    # P02
    s = slide_base(prs, "为什么今天要谈开发模式", "导入", 2)
    icon_building(s, 0.70, 1.34, BLUE)
    text_box(s, "过去常问", 1.34, 1.42, 3.0, 0.34, BODY_PT, BLUE_DARK, True)
    icon_agent(s, 7.06, 1.34, GREEN)
    text_box(s, "现在更应关注", 7.68, 1.42, 3.2, 0.34, BODY_PT, GREEN, True)
    left = ["能否更快写代码", "谁负责完成任务", "何时交付功能"]
    right = ["如何形成正确结果", "谁约束并审查结果", "如何持续验证价值"]
    for i, (lv, rv) in enumerate(zip(left, right)):
        y = 2.12 + i * 0.88
        box(s, 0.70, y, 4.68, 0.60, BLUE_LIGHT, BLUE_LIGHT)
        text_box(s, lv, 0.94, y + 0.13, 4.05, 0.32, BODY_PT, TEXT)
        chevron(s, 6.04, y + 0.20, 0.30, 0.20, BLUE)
        box(s, 7.05, y, 5.40, 0.60, GREEN_LIGHT, GREEN_LIGHT)
        text_box(s, rv, 7.30, y + 0.13, 4.78, 0.32, BODY_PT, TEXT, True)
    box(s, 0.70, 5.36, 11.75, 0.75, WHITE, BLUE)
    text_box(s, "Agent 改变的不只是编码速度", 0.98, 5.57, 11.20, 0.34,
             BODY_PT, BLUE_DARK, True, PP_ALIGN.CENTER)

    # P03
    s = slide_base(prs, "本课程回答什么", "导入", 3)
    header_card(s, 0.72, 1.42, 5.62, 4.58, "聚焦内容",
                ["开发模式如何变化", "生命周期如何重组", "管理责任如何变化", "价值风险如何平衡"],
                BLUE, WHITE, icon_target)
    header_card(s, 6.92, 1.42, 5.62, 4.58, "本次不展开",
                ["工具安装与命令", "团队导入方案", "产品选型比较", "现场实操训练"],
                ORANGE, WHITE, icon_shield)

    # P04
    s = slide_base(prs, "瀑布/V 模型：用阶段控制复杂度", "模式演进", 4)
    labels = ["需求", "设计", "开发", "测试", "交付"]
    icons = [icon_document, icon_target, icon_agent, icon_check, icon_building]
    for i, (label, ic) in enumerate(zip(labels, icons)):
        x = 0.92 + i * 2.32
        ic(s, x + 0.35, 1.62, BLUE)
        box(s, x, 2.28, 1.20, 0.58, BLUE_LIGHT, BLUE_LIGHT)
        text_box(s, label, x + 0.08, 2.41, 1.04, 0.30, BODY_PT, BLUE_DARK, True, PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            connector(s, x + 1.22, 2.57, x + 2.12, 2.57, BLUE, 2)
            chevron(s, x + 1.90, 2.45, 0.20, 0.23, BLUE)
    connector(s, 1.52, 3.32, 6.10, 4.50, GRAY_TEXT, 1)
    connector(s, 6.10, 4.50, 10.80, 3.32, GRAY_TEXT, 1)
    pill(s, "需求对应验收", 0.90, 4.98, 2.35, BLUE, BLUE_LIGHT)
    pill(s, "设计对应测试", 9.10, 4.98, 2.35, BLUE, BLUE_LIGHT)
    small_stat(s, 0.90, 6.02, 5.26, "价值", "范围清晰 · 责任可追踪", GREEN)
    small_stat(s, 6.66, 6.02, 5.26, "限制", "反馈较晚 · 变更成本高", ORANGE)

    # P05
    s = slide_base(prs, "敏捷与 DevOps：用反馈缩短周期", "模式演进", 5)
    nodes = [("计划", 6.02, 1.44), ("构建", 8.25, 2.43), ("测试", 7.58, 4.72),
             ("发布", 4.34, 4.72), ("反馈", 3.68, 2.43)]
    for text, x, y in nodes:
        circle(s, x, y, 0.76, BLUE_LIGHT if text != "反馈" else GREEN_LIGHT,
               BLUE if text != "反馈" else GREEN)
        text_box(s, text, x + 0.04, y + 0.22, 0.68, 0.30, BODY_PT,
                 BLUE_DARK if text != "反馈" else GREEN, True, PP_ALIGN.CENTER)
    paths = [(6.78, 1.82, 8.24, 2.55), (8.45, 3.18, 7.96, 4.61),
             (7.55, 5.10, 5.12, 5.10), (4.33, 4.62, 3.94, 3.18),
             (4.34, 2.43, 5.98, 1.83)]
    for x1, y1, x2, y2 in paths:
        connector(s, x1, y1, x2, y2, BLUE, 2)
    header_card(s, 0.72, 1.56, 2.42, 3.52, "核心变化",
                ["小批次交付", "持续测试验证", "运行反馈回流"], BLUE, WHITE, icon_network)
    box(s, 9.55, 1.78, 2.72, 1.30, GREEN_LIGHT, GREEN_LIGHT)
    text_box(s, "价值", 9.82, 2.01, 2.18, 0.28, 16, GREEN, True)
    text_box(s, "更快验证价值", 9.82, 2.42, 2.18, 0.34, BODY_PT, TEXT, True)
    box(s, 9.55, 3.58, 2.72, 1.30, ORANGE_LIGHT, ORANGE_LIGHT)
    text_box(s, "挑战", 9.82, 3.81, 2.18, 0.28, 16, ORANGE, True)
    text_box(s, "持续保持质量", 9.82, 4.22, 2.20, 0.34, BODY_PT, TEXT, True)

    # P06
    s = slide_base(prs, "Agent 时代继承什么", "模式演进", 6)
    stacks = [
        (4.07, 4.62, 5.18, 0.72, BLUE_LIGHT, BLUE, "瀑布 / V：基线、审查与责任"),
        (3.48, 3.72, 6.36, 0.72, RGBColor(0xD9, 0xEA, 0xF7), BLUE, "敏捷 / DevOps：迭代与持续验证"),
        (2.88, 2.82, 7.56, 0.72, GREEN_LIGHT, GREEN, "AI Agent：任务委派与密集反馈"),
    ]
    for x, y, w, h, fill, accent, label in stacks:
        box(s, x, y, w, h, fill, fill)
        fill_bar(s, x, y, 0.08, h, accent)
        text_box(s, label, x + 0.28, y + 0.20, w - 0.50, 0.30, BODY_PT, BLACK, True, PP_ALIGN.CENTER)
    icon_document(s, 1.12, 4.77, BLUE)
    icon_network(s, 1.12, 3.87, BLUE)
    icon_agent(s, 1.12, 2.97, GREEN)
    pill(s, "可靠交付", 5.34, 5.87, 2.60, BLUE, BLUE_LIGHT)
    box(s, 2.18, 1.55, 8.95, 0.62, WHITE, BLUE)
    text_box(s, "不是替代，而是重组", 2.38, 1.71, 8.56, 0.30, BODY_PT, BLUE_DARK, True, PP_ALIGN.CENTER)

    # P07
    s = slide_base(prs, "什么是本课程所说的 AI Agent", "Agent 定义", 7)
    stages = [
        ("问答辅助", "回答问题", BLUE_LIGHT, BLUE),
        ("内容生成", "生成片段", BLUE_LIGHT, BLUE),
        ("工具执行", "读取与检查", BLUE_LIGHT, BLUE),
        ("Agent 协作", "修改并验证", GREEN_LIGHT, GREEN),
    ]
    for i, (name, value, fill, accent) in enumerate(stages):
        x = 0.83 + i * 3.01
        y = 4.55 - i * 0.63
        box(s, x, y, 2.37, 1.18, fill, fill)
        fill_bar(s, x, y, 0.07, 1.18, accent)
        text_box(s, name, x + 0.20, y + 0.20, 2.00, 0.29, BODY_PT, accent, True)
        text_box(s, value, x + 0.20, y + 0.64, 2.00, 0.30, BODY_PT, TEXT)
        if i < 3:
            connector(s, x + 2.38, y + 0.62, x + 2.96, y - 0.03, BLUE, 1.5)
    box(s, 1.42, 6.11, 10.48, 0.64, WHITE, BLUE)
    text_box(s, "上下文  +  授权  +  执行  +  验证", 1.64, 6.29, 10.04, 0.28,
             BODY_PT, BLUE_DARK, True, PP_ALIGN.CENTER)

    # P08
    s = slide_base(prs, "人与 Agent 的职责边界", "Agent 定义", 8)
    header_card(s, 0.72, 1.40, 5.54, 4.72, "Agent 可辅助",
                ["资料整理", "方案比较", "任务执行", "测试运行", "结果汇报"],
                BLUE, BLUE_LIGHT, icon_agent)
    header_card(s, 7.05, 1.40, 5.54, 4.72, "人必须负责",
                ["业务目标", "范围取舍", "关键批准", "证据审查", "最终责任"],
                GREEN, GREEN_LIGHT, icon_building)
    chevron(s, 6.48, 3.59, 0.35, 0.30, BLUE)
    text_box(s, "责任不能空缺", 4.56, 6.42, 4.30, 0.35, BODY_PT, BLUE_DARK, True, PP_ALIGN.CENTER)

    # P09
    s = slide_base(prs, "已经出现的实践框架", "Agent 定义", 9)
    items = [
        ("Superpowers", ["澄清", "计划", "测试", "审查"], "用流程约束执行", BLUE),
        ("Spec Kit", ["规格", "计划", "任务", "实现"], "规格驱动实现", GREEN),
        ("OpenSpec", ["提案", "规格", "设计", "任务"], "变更保持可追踪", BLUE),
    ]
    for i, (name, flow, takeaway, accent) in enumerate(items):
        x = 0.72 + i * 4.03
        box(s, x, 1.54, 3.62, 4.36, WHITE, GRAY)
        text_box(s, name, x + 0.25, 1.82, 3.14, 0.33, BODY_PT, accent, True)
        for j, word in enumerate(flow):
            y = 2.44 + j * 0.55
            pill(s, word, x + 0.30, y, 1.12, accent,
                 GREEN_LIGHT if accent == GREEN else BLUE_LIGHT)
            if j < 3:
                chevron(s, x + 1.58, y + 0.10, 0.17, 0.18, accent)
                text_box(s, flow[j + 1], x + 1.94, y + 0.05, 1.28, 0.25, 15, TEXT)
                break
        # Compact second workflow line.
        text_box(s, f"{flow[2]}  ->  {flow[3]}", x + 0.32, 3.60, 2.94, 0.33, BODY_PT, TEXT)
        fill_bar(s, x + 0.25, 4.34, 3.10, 0.03, accent)
        text_box(s, takeaway, x + 0.26, 4.67, 3.10, 0.41, BODY_PT, BLACK, True, PP_ALIGN.CENTER)
    box(s, 2.16, 6.15, 9.00, 0.52, WHITE, BLUE)
    text_box(s, "从即时生成，走向工程协作", 2.32, 6.27, 8.68, 0.28,
             BODY_PT, BLUE_DARK, True, PP_ALIGN.CENTER)
    text_box(s, "公开仓库：obra/superpowers · github/spec-kit · Fission-AI/OpenSpec",
             0.73, 6.78, 8.40, 0.20, META_PT, GRAY_TEXT)

    # P10
    s = slide_base(prs, "全生命周期：变化发生在哪里", "生命周期", 10)
    phase_names = ["构思需求", "规划设计", "编码", "测试", "发布", "运维", "反馈"]
    phase_icons = [icon_target, icon_document, icon_agent, icon_check, icon_building, icon_shield, icon_network]
    for i, (name, ic) in enumerate(zip(phase_names, phase_icons)):
        x = 0.58 + i * 1.78
        ic(s, x + 0.34, 2.25, BLUE if name not in ("测试", "反馈") else GREEN)
        text_box(s, name, x + 0.02, 2.90, 1.14, 0.34, BODY_PT, TEXT, True, PP_ALIGN.CENTER)
        pill(s, "+ Agent", x + 0.12, 3.65, 1.04, GREEN, GREEN_LIGHT)
        if i < 6:
            connector(s, x + 1.25, 2.55, x + 1.65, 2.55, BLUE, 1.5)
            chevron(s, x + 1.52, 2.45, 0.14, 0.19, BLUE)
    box(s, 2.08, 5.27, 9.20, 0.80, WHITE, BLUE)
    text_box(s, "Agent 参与的不只是编码", 2.30, 5.51, 8.76, 0.32,
             BODY_PT, BLUE_DARK, True, PP_ALIGN.CENTER)

    # P11
    s = slide_base(prs, "构思与需求：从记录到辅助澄清", "生命周期", 11)
    cols = [
        ("传统方式", ["会议访谈", "人工归纳", "文档评审"], BLUE, icon_document),
        ("Agent 可辅助", ["连续追问", "汇总分歧", "检查遗漏"], GREEN, icon_agent),
        ("管理关注", ["目标确认", "范围取舍", "假设标识"], ORANGE, icon_shield),
    ]
    for i, (head, items, accent, ic) in enumerate(cols):
        header_card(s, 0.73 + i * 4.03, 1.56, 3.56, 3.90, head, items, accent, WHITE, ic)
    box(s, 2.04, 5.94, 9.27, 0.64, WHITE, BLUE)
    text_box(s, "澄清可加速，取舍仍在人", 2.24, 6.12, 8.86, 0.30,
             BODY_PT, BLUE_DARK, True, PP_ALIGN.CENTER)

    # P12
    s = slide_base(prs, "规划与设计：从交接到上下文协作", "生命周期", 12)
    flow = [("需求", "分层交接"), ("设计", "检索约束"), ("任务", "核对一致"), ("评审", "批准承诺")]
    for i, (head, sub) in enumerate(flow):
        x = 0.80 + i * 3.05
        box(s, x, 2.02, 2.46, 1.32, BLUE_LIGHT if i < 3 else GREEN_LIGHT,
            BLUE_LIGHT if i < 3 else GREEN_LIGHT)
        text_box(s, head, x + 0.16, 2.28, 2.14, 0.30, BODY_PT,
                 GREEN if i == 3 else BLUE_DARK, True, PP_ALIGN.CENTER)
        text_box(s, sub, x + 0.16, 2.78, 2.14, 0.28, BODY_PT, TEXT, False, PP_ALIGN.CENTER)
        if i < 3:
            chevron(s, x + 2.62, 2.55, 0.25, 0.24, BLUE)
    small_stat(s, 0.80, 4.38, 3.56, "架构承诺", "由人批准", GREEN)
    small_stat(s, 4.87, 4.38, 3.56, "成本影响", "由人判断", BLUE)
    small_stat(s, 8.94, 4.38, 3.56, "安全边界", "由人负责", ORANGE)
    text_box(s, "方案可推演，承诺需审批", 3.76, 6.20, 5.80, 0.35,
             BODY_PT, BLUE_DARK, True, PP_ALIGN.CENTER)

    # P13
    s = slide_base(prs, "编码与测试：从执行到受控委派", "生命周期", 13)
    loop = [("读取", 5.90, 1.65), ("修改", 8.27, 2.60), ("测试", 7.45, 4.55),
            ("报告", 4.35, 4.55), ("审查", 3.63, 2.60)]
    for word, x, y in loop:
        accent = GREEN if word in ("测试", "报告") else BLUE
        circle(s, x, y, 0.78, GREEN_LIGHT if accent == GREEN else BLUE_LIGHT, accent)
        text_box(s, word, x + 0.06, y + 0.22, 0.66, 0.31, BODY_PT, accent, True, PP_ALIGN.CENTER)
    for x1, y1, x2, y2 in [(6.68, 2.00, 8.20, 2.72), (8.46, 3.37, 7.85, 4.48),
                           (7.43, 4.96, 5.15, 4.96), (4.30, 4.50, 3.91, 3.34),
                           (4.36, 2.62, 5.87, 2.02)]:
        connector(s, x1, y1, x2, y2, BLUE, 1.6)
    header_card(s, 0.66, 1.68, 2.44, 3.80, "管理关注",
                ["不扩大范围", "测试真执行", "高风险复核"], ORANGE, WHITE, icon_shield)
    box(s, 9.82, 2.12, 2.30, 2.34, GREEN_LIGHT, GREEN_LIGHT)
    icon_check(s, 10.71, 2.42, GREEN)
    text_box(s, "执行更密集", 10.02, 3.20, 1.90, 0.32, BODY_PT, GREEN, True, PP_ALIGN.CENTER)
    text_box(s, "证据更重要", 10.02, 3.70, 1.90, 0.32, BODY_PT, BLACK, True, PP_ALIGN.CENTER)

    # P14
    s = slide_base(prs, "发布与运维：从执行到辅助诊断", "生命周期", 14)
    gate = ["准备", "审批", "上线", "监控", "诊断"]
    for i, label in enumerate(gate):
        x = 0.94 + i * 2.35
        fill = GREEN_LIGHT if label in ("监控", "诊断") else BLUE_LIGHT
        accent = GREEN if label in ("监控", "诊断") else BLUE
        box(s, x, 1.68, 1.64, 0.66, fill, fill)
        text_box(s, label, x + 0.14, 1.87, 1.36, 0.29, BODY_PT, accent, True, PP_ALIGN.CENTER)
        if i < 4:
            chevron(s, x + 1.84, 1.90, 0.21, 0.21, BLUE)
    header_card(s, 0.94, 3.02, 5.25, 2.35, "Agent 可辅助",
                ["整理发布内容", "归纳日志信息", "提出定位方向"], GREEN, WHITE, icon_check)
    header_card(s, 7.09, 3.02, 5.25, 2.35, "不可失控委派",
                ["上线批准", "敏感权限", "事故决策"], ORANGE, WHITE, icon_shield)
    text_box(s, "诊断可加速，权限需守住", 3.66, 6.02, 6.03, 0.36,
             BODY_PT, BLUE_DARK, True, PP_ALIGN.CENTER)

    # P15
    s = slide_base(prs, "迭代反馈：从复盘到快速归纳", "生命周期", 15)
    header_card(s, 0.84, 1.72, 4.86, 3.34, "可加速工作",
                ["汇总用户反馈", "整理缺陷模式", "同步变更记录"], GREEN, GREEN_LIGHT, icon_network)
    header_card(s, 7.63, 1.72, 4.86, 3.34, "人的判断",
                ["识别真实价值", "决定优先顺序", "确认下轮范围"], BLUE, BLUE_LIGHT, icon_building)
    connector(s, 5.74, 3.38, 7.40, 3.38, GREEN, 2)
    chevron(s, 7.18, 3.26, 0.19, 0.22, GREEN)
    connector(s, 7.64, 5.62, 1.40, 5.62, BLUE, 1.5)
    chevron(s, 1.25, 5.50, 0.18, 0.22, BLUE)
    pill(s, "回流下一轮需求", 4.92, 5.38, 3.18, BLUE, BLUE_LIGHT)
    text_box(s, "信息可归纳，优先级在人", 3.50, 6.35, 6.42, 0.32,
             BODY_PT, BLUE_DARK, True, PP_ALIGN.CENTER)

    # P16
    s = slide_base(prs, "Agent 带来的三类价值", "价值与风险", 16)
    values = [
        ("信息处理更快", "快速阅读与归纳", icon_document),
        ("环节衔接更顺", "同步需求与产物", icon_network),
        ("验证反馈更密", "修改后立即检查", icon_check),
    ]
    for i, (head, detail, ic) in enumerate(values):
        x = 0.79 + i * 4.12
        box(s, x, 1.70, 3.60, 3.46, GREEN_LIGHT, GREEN_LIGHT)
        ic(s, x + 1.56, 2.08, GREEN)
        text_box(s, head, x + 0.20, 2.87, 3.20, 0.33, BODY_PT, GREEN, True, PP_ALIGN.CENTER)
        text_box(s, detail, x + 0.20, 3.62, 3.20, 0.32, BODY_PT, TEXT, False, PP_ALIGN.CENTER)
        if i < 2:
            connector(s, x + 3.62, 3.39, x + 4.04, 3.39, BLUE, 1.5)
    box(s, 2.06, 5.83, 9.22, 0.68, WHITE, GREEN)
    text_box(s, "价值来自更短的反馈距离", 2.30, 6.03, 8.74, 0.30,
             BODY_PT, GREEN, True, PP_ALIGN.CENTER)

    # P17
    s = slide_base(prs, "速度不会自动变成质量", "价值与风险", 17)
    risks = [
        ("错误落实更快", "假设进入实现"),
        ("完成错觉更强", "生成不等于通过"),
        ("权限影响更大", "越界修改数据"),
        ("责任容易模糊", "输出无人负责"),
    ]
    for i, (head, detail) in enumerate(risks):
        x = 0.68 + i * 3.14
        box(s, x, 1.82, 2.68, 3.08, ORANGE_LIGHT, ORANGE_LIGHT)
        icon_shield(s, x + 1.10, 2.14, ORANGE)
        text_box(s, head, x + 0.14, 2.98, 2.40, 0.34, BODY_PT, ORANGE, True, PP_ALIGN.CENTER)
        text_box(s, detail, x + 0.14, 3.70, 2.40, 0.34, BODY_PT, TEXT, False, PP_ALIGN.CENTER)
    fill_bar(s, 1.40, 5.43, 10.54, 0.06, ORANGE)
    chevron(s, 11.84, 5.34, 0.28, 0.22, ORANGE)
    text_box(s, "加速也会放大错误", 4.58, 5.83, 4.36, 0.38,
             BODY_PT, ORANGE, True, PP_ALIGN.CENTER)

    # P18
    s = slide_base(prs, "新模式仍需要的控制原则", "价值与风险", 18)
    principles = [("01", "目标明确", BLUE), ("02", "上下文充分", BLUE),
                  ("03", "授权有界", ORANGE), ("04", "验证有据", GREEN),
                  ("05", "责任有人", BLUE)]
    for i, (num, label, accent) in enumerate(principles):
        x = 0.77 + i * 2.49
        circle(s, x + 0.74, 2.04, 0.58, WHITE, accent)
        text_box(s, num, x + 0.78, 2.20, 0.50, 0.23, 14, accent, True, PP_ALIGN.CENTER)
        box(s, x, 2.92, 2.05, 0.72,
            ORANGE_LIGHT if accent == ORANGE else GREEN_LIGHT if accent == GREEN else BLUE_LIGHT,
            WHITE)
        text_box(s, label, x + 0.10, 3.14, 1.85, 0.30, BODY_PT, accent, True, PP_ALIGN.CENTER)
        if i < 4:
            connector(s, x + 2.06, 3.28, x + 2.40, 3.28, BLUE, 1.5)
            chevron(s, x + 2.29, 3.17, 0.12, 0.19, BLUE)
    box(s, 2.34, 5.16, 8.67, 0.88, WHITE, BLUE)
    text_box(s, "可控协作，才是可靠加速", 2.56, 5.45, 8.24, 0.32,
             BODY_PT, BLUE_DARK, True, PP_ALIGN.CENTER)

    # P19
    s = slide_base(prs, "三种模式不是三选一", "总结", 19)
    layers = [
        ("瀑布 / V 模型", "基线与验证责任", BLUE),
        ("敏捷 / DevOps", "迭代与持续反馈", GREEN),
        ("AI Agent", "细粒度受控执行", BLUE_DARK),
    ]
    for i, (name, value, accent) in enumerate(layers):
        y = 1.66 + i * 1.18
        box(s, 1.12, y, 4.12, 0.83, BLUE_LIGHT if accent != GREEN else GREEN_LIGHT, WHITE)
        text_box(s, name, 1.34, y + 0.25, 3.70, 0.30, BODY_PT, accent, True)
        chevron(s, 5.74, y + 0.29, 0.30, 0.22, BLUE)
        box(s, 6.48, y, 4.76, 0.83, WHITE, GRAY)
        text_box(s, value, 6.70, y + 0.25, 4.32, 0.30, BODY_PT, TEXT, True)
    box(s, 2.72, 5.62, 7.88, 0.78, GREEN_LIGHT, GREEN)
    icon_building(s, 3.24, 5.77, GREEN)
    text_box(s, "组合能力，形成可靠交付", 4.01, 5.84, 5.98, 0.35,
             BODY_PT, GREEN, True, PP_ALIGN.CENTER)

    # P20
    s = slide_base(prs, "管理者需要回答的五个问题", "总结", 20)
    questions = ["目标是否明确？", "上下文是否充分？", "授权是否适当？", "结果是否验证？", "责任是否清晰？"]
    icons = [icon_target, icon_document, icon_shield, icon_check, icon_building]
    accents = [BLUE, BLUE, ORANGE, GREEN, BLUE]
    for i, (question, ic, accent) in enumerate(zip(questions, icons, accents)):
        x = 0.68 + (i % 3) * 4.08
        y = 1.62 + (i // 3) * 1.83
        if i >= 3:
            x += 2.04
        box(s, x, y, 3.56, 1.22, WHITE, GRAY)
        text_box(s, f"0{i + 1}", x + 0.14, y + 0.18, 0.40, 0.28, BODY_PT, accent, True)
        ic(s, x + 0.57, y + 0.16, accent)
        text_box(s, question, x + 1.10, y + 0.46, 2.38, 0.32, BODY_PT, BLACK, True)
    box(s, 2.62, 5.47, 8.11, 0.65, BLUE_LIGHT, BLUE_LIGHT)
    text_box(s, "评估 Agent 协作成熟度的共同语言", 2.82, 5.66, 7.72, 0.30,
             BODY_PT, BLUE_DARK, True, PP_ALIGN.CENTER)

    # P21
    s = slide_base(prs, "结语", "总结", 21)
    icon_agent(s, 11.26, 1.50, BLUE)
    icon_building(s, 11.82, 1.50, GREEN)
    box(s, 0.85, 1.56, 9.76, 2.10, BLUE_LIGHT, BLUE_LIGHT)
    text_box(s, "Agent 加速软件工作", 1.20, 2.00, 9.04, 0.36, BODY_PT, BLUE_DARK, True, PP_ALIGN.CENTER)
    text_box(s, "人定义方向并对结果负责", 1.20, 2.67, 9.04, 0.36, BODY_PT, BLACK, True, PP_ALIGN.CENTER)
    closing = [("保留工程纪律", BLUE), ("提升反馈密度", GREEN), ("守住责任边界", ORANGE)]
    for i, (value, accent) in enumerate(closing):
        x = 1.20 + i * 3.78
        box(s, x, 4.78, 3.18, 0.76,
            GREEN_LIGHT if accent == GREEN else ORANGE_LIGHT if accent == ORANGE else BLUE_LIGHT,
            WHITE)
        text_box(s, value, x + 0.14, 5.02, 2.90, 0.30, BODY_PT, accent, True, PP_ALIGN.CENTER)

    # P22
    s = slide_base(prs, "参考资料", "参考", 22)
    refs = [
        ("Superpowers", "软件开发方法与技能框架", "https://github.com/obra/superpowers", icon_agent),
        ("Spec Kit", "规格驱动开发工具包", "https://github.com/github/spec-kit", icon_document),
        ("OpenSpec", "轻量规格与变更组织", "https://github.com/Fission-AI/OpenSpec", icon_network),
    ]
    for i, (name, description, url, ic) in enumerate(refs):
        y = 1.56 + i * 1.46
        box(s, 0.92, y, 11.48, 1.02, WHITE, GRAY)
        ic(s, 1.18, y + 0.27, BLUE)
        text_box(s, name, 1.91, y + 0.18, 2.55, 0.29, BODY_PT, BLUE_DARK, True)
        text_box(s, description, 4.62, y + 0.18, 3.56, 0.30, BODY_PT, TEXT)
        link = text_box(s, "查看公开仓库", 9.52, y + 0.18, 2.23, 0.30, BODY_PT, BLUE, True, PP_ALIGN.RIGHT)
        for p in link.text_frame.paragraphs:
            for run in p.runs:
                run.hyperlink.address = url
        fill_bar(s, 9.62, y + 0.58, 2.09, 0.025, BLUE)
    text_box(s, "仅作为工程化实践示例，不构成选型建议", 0.94, 6.22, 11.44, 0.30,
             BODY_PT, GRAY_TEXT, False, PP_ALIGN.CENTER)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    output = build_deck()
    print(output)
