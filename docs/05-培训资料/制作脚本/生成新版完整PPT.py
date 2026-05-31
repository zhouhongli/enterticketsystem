# -*- coding: utf-8 -*-
"""Generate the redesigned full training presentation based on the v1.0 page script."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parent.parent / "08-传统软件开发模式与AI-Agent时代的软件开发模式-新版完整母版.pptx"

FONT = "Microsoft YaHei"
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x0E, 0x18, 0x25)
TEXT = RGBColor(0x34, 0x40, 0x52)
MUTED = RGBColor(0x6F, 0x7C, 0x8E)
LINE = RGBColor(0xE4, 0xEA, 0xF0)
GRID = RGBColor(0xF3, 0xF6, 0xF9)
BLUE = RGBColor(0x07, 0x56, 0xC7)
BLUE_DARK = RGBColor(0x04, 0x3D, 0x93)
BLUE_PALE = RGBColor(0xF1, 0xF6, 0xFC)
GREEN = RGBColor(0x0C, 0xA4, 0x69)
GREEN_PALE = RGBColor(0xED, 0xF8, 0xF3)
ORANGE = RGBColor(0xF3, 0x83, 0x21)
ORANGE_PALE = RGBColor(0xFE, 0xF4, 0xE9)

SW = 13.333
SH = 7.5
TOTAL = 39


def I(n):
    return Inches(n)


def east_asian(run):
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    if ea is None:
        ea = OxmlElement("a:ea")
        rpr.append(ea)
    ea.set("typeface", FONT)


def run_style(run, size, color, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    east_asian(run)


def text(slide, value, x, y, w, h, size=18, color=TEXT, bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE, spacing=1.04):
    shape = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = I(0)
    tf.margin_right = I(0)
    tf.margin_top = I(0)
    tf.margin_bottom = I(0)
    tf.vertical_anchor = valign
    for idx, line_value in enumerate(value.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line_value
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(0)
        for run in p.runs:
            run_style(run, size, color, bold)
    return shape


def rect(slide, x, y, w, h, fill=WHITE, line_color=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, I(x), I(y), I(w), I(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    return shape


def rule(slide, x, y, w, color=BLUE, thickness=0.035):
    return rect(slide, x, y, w, thickness, color)


def line(slide, x1, y1, x2, y2, color=LINE, width=1.2):
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    return shape


def dot(slide, x, y, d, fill, outline=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, I(x), I(y), I(d), I(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if outline is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = outline
        shape.line.width = Pt(1.2)
    return shape


def tag(slide, value, x, y, w, color=BLUE, fill=BLUE_PALE):
    rect(slide, x, y, w, 0.34, fill, None, True)
    text(slide, value, x + 0.03, y + 0.065, w - 0.06, 0.22, 11, color, True, PP_ALIGN.CENTER)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide


def chrome(slide, section, number):
    text(slide, section, 0.62, 0.29, 2.8, 0.19, 10, BLUE, True)
    rule(slide, 0.62, 0.57, 12.07, LINE, 0.018)
    rule(slide, 0.62, 7.06, 12.07, LINE, 0.018)
    text(slide, "传统软件开发模式与 AI Agent 时代的软件开发模式", 0.62, 7.16, 5.2, 0.17, 9, MUTED)
    text(slide, f"{number:02d}  /  {TOTAL:02d}", 11.84, 7.15, 0.84, 0.18, 10, MUTED, True, PP_ALIGN.RIGHT)


def page(prs, section, number, title_value, subtitle=None, key=False):
    slide = new_slide(prs)
    chrome(slide, section, number)
    size = 34 if key else 28
    text(slide, title_value, 0.62, 0.82, 11.6, 0.48 if not key else 0.56, size, BLACK, True)
    if subtitle:
        text(slide, subtitle, 0.64, 1.43 if not key else 1.56, 11.2, 0.27, 15, MUTED)
    return slide


def bullet(slide, value, x, y, w, accent=BLUE, size=17):
    dot(slide, x, y + 0.11, 0.075, accent)
    text(slide, value, x + 0.18, y, w - 0.18, 0.30, size, TEXT)


def card(slide, x, y, w, h, header, lines, accent=BLUE, fill=WHITE,
         body_size=16):
    rect(slide, x, y, w, h, fill, LINE, True)
    rule(slide, x + 0.18, y + 0.18, 0.48, accent, 0.045)
    text(slide, header, x + 0.18, y + 0.36, w - 0.36, 0.33, 18, BLACK, True)
    for idx, item in enumerate(lines):
        bullet(slide, item, x + 0.20, y + 0.86 + idx * 0.39, w - 0.38, accent, body_size)


def statement(slide, value, x, y, w, accent=BLUE_DARK, size=22):
    rule(slide, x, y, 0.65, accent, 0.045)
    text(slide, value, x, y + 0.18, w, 0.36, size, accent, True)


def numbered(slide, value, x, y, accent=BLUE, d=0.32):
    dot(slide, x, y, d, WHITE, accent)
    text(slide, value, x, y + 0.08, d, 0.17, 10, accent, True, PP_ALIGN.CENTER)


def flow_nodes(slide, labels, x, y, width, accent=BLUE, size=14):
    count = len(labels)
    node_w = (width - (count - 1) * 0.18) / count
    for idx, item in enumerate(labels):
        nx = x + idx * (node_w + 0.18)
        rect(slide, nx, y, node_w, 0.56, BLUE_PALE if accent == BLUE else GREEN_PALE, None, True)
        text(slide, item, nx + 0.03, y + 0.13, node_w - 0.06, 0.28, size, accent, True, PP_ALIGN.CENTER)
        if idx < count - 1:
            text(slide, ">", nx + node_w, y + 0.16, 0.18, 0.24, 13, MUTED, True, PP_ALIGN.CENTER)


def doc_sheet(slide, x, y, w, h, label, entries, accent=BLUE):
    rect(slide, x, y, w, h, WHITE, LINE, True)
    rect(slide, x, y, w, 0.40, BLUE_PALE if accent == BLUE else GREEN_PALE, None, True)
    text(slide, label, x + 0.14, y + 0.10, w - 0.28, 0.19, 12, accent, True)
    yy = y + 0.58
    for code, value in entries:
        if code:
            tag(slide, code, x + 0.15, yy + 0.02, min(0.90, 0.13 + len(code) * 0.075), accent,
                BLUE_PALE if accent == BLUE else GREEN_PALE)
            tx = x + 1.16
        else:
            tx = x + 0.18
        text(slide, value, tx, yy + 0.03, w - (tx - x) - 0.12, 0.30, 13, TEXT)
        line(slide, x + 0.15, yy + 0.44, x + w - 0.15, yy + 0.44, LINE, 0.7)
        yy += 0.54


def table(slide, x, y, widths, headers, rows, row_h=0.50, size=13,
          accent_col=None):
    total_w = sum(widths)
    rect(slide, x, y, total_w, row_h, BLUE_DARK, None, False)
    xx = x
    for width, header in zip(widths, headers):
        text(slide, header, xx + 0.10, y + 0.12, width - 0.20, row_h - 0.18, size, WHITE, True)
        xx += width
    for row_idx, row in enumerate(rows):
        yy = y + (row_idx + 1) * row_h
        fill = WHITE if row_idx % 2 == 0 else BLUE_PALE
        rect(slide, x, yy, total_w, row_h, fill, None, False)
        xx = x
        for col_idx, (width, value) in enumerate(zip(widths, row)):
            color = accent_col if col_idx == len(row) - 1 and accent_col else TEXT
            bold = col_idx == 0
            text(slide, value, xx + 0.10, yy + 0.10, width - 0.20, row_h - 0.15, size, color, bold)
            xx += width
    rect(slide, x, y, total_w, row_h * (len(rows) + 1), WHITE, LINE, False).fill.background()


def panorama(slide, roles, processes, artifacts, accent=BLUE, case_note=""):
    text(slide, "参与角色", 0.68, 2.01, 0.92, 0.24, 12, MUTED, True)
    text(slide, "工作流程", 0.68, 3.23, 0.92, 0.24, 12, MUTED, True)
    text(slide, "形成产物", 0.68, 4.57, 0.92, 0.24, 12, MUTED, True)
    line(slide, 1.77, 1.94, 1.77, 5.30, LINE, 1)
    flow_nodes(slide, roles, 1.98, 1.84, 10.45, BLUE, 12)
    flow_nodes(slide, processes, 1.98, 3.05, 10.45, accent, 12)
    flow_nodes(slide, artifacts, 1.98, 4.38, 10.45, GREEN, 12)
    if case_note:
        tag(slide, case_note, 1.98, 5.60, 5.55, BLUE, BLUE_PALE)


def gate_chain(slide, items, conclusion):
    start = 0.82
    width = 11.70
    line(slide, start + 0.28, 3.25, start + width - 0.28, 3.25, LINE, 1.5)
    gap = width / len(items)
    for idx, (label, owner, color) in enumerate(items):
        x = start + idx * gap
        dot(slide, x + 0.27, 3.03, 0.44, color)
        text(slide, str(idx + 1).zfill(2), x + 0.27, 3.17, 0.44, 0.14, 9, WHITE, True, PP_ALIGN.CENTER)
        text(slide, label, x, 2.36, gap - 0.12, 0.31, 16, color, True, PP_ALIGN.CENTER)
        text(slide, owner, x - 0.03, 3.78, gap, 0.48, 13, TEXT, False, PP_ALIGN.CENTER)
    rect(slide, 0.82, 5.28, 11.64, 0.66, BLUE_DARK, None, True)
    text(slide, conclusion, 0.94, 5.46, 11.40, 0.26, 20, WHITE, True, PP_ALIGN.CENTER)


def framework(slide, flow, observations, source, accent):
    text(slide, "公开工作流示意", 0.72, 2.03, 1.54, 0.22, 12, MUTED, True)
    flow_nodes(slide, flow, 0.72, 2.43, 11.85, accent, 14)
    text(slide, "课程解读", 0.72, 3.55, 1.54, 0.22, 12, MUTED, True)
    for idx, item in enumerate(observations):
        rect(slide, 0.72 + idx * 3.93, 3.98, 3.52, 0.92, WHITE, LINE, True)
        rule(slide, 0.90 + idx * 3.93, 4.15, 0.42, accent, 0.04)
        text(slide, item, 0.90 + idx * 3.93, 4.39, 3.14, 0.28, 16, TEXT, True)
    tag(slide, "实例说明，不代表统一标准", 0.72, 5.60, 2.83, ORANGE, ORANGE_PALE)
    text(slide, source, 0.72, 6.25, 11.80, 0.20, 11, MUTED)


def cover(prs):
    slide = new_slide(prs)
    chrome(slide, "技术管理培训  /  新版完整母版", 1)
    text(slide, "传统软件开发模式", 0.62, 1.16, 6.85, 0.55, 38, BLACK, True)
    text(slide, "与 AI Agent 时代的", 0.62, 1.82, 7.30, 0.55, 38, BLACK, True)
    text(slide, "软件开发模式", 0.62, 2.48, 6.40, 0.55, 38, BLUE, True)
    text(slide, "面向产品、项目与研发管理人员", 0.64, 3.38, 5.80, 0.30, 17, MUTED)
    text(slide, "PROCESS  /  FEEDBACK  /  CONTROLLED EXECUTION", 8.14, 1.22, 4.40, 0.22, 9, MUTED, True)
    stages = [("01", "阶段控制", "瀑布 / V", BLUE), ("02", "持续反馈", "敏捷 / DevOps", BLUE),
              ("03", "受控协作", "AI Agent", GREEN)]
    line(slide, 8.18, 1.87, 12.27, 1.87, BLUE, 2)
    for idx, (no, head, sub, accent) in enumerate(stages):
        y = 2.23 + idx * 1.12
        numbered(slide, no, 8.18, y + 0.03, accent, 0.35)
        text(slide, head, 8.78, y, 2.36, 0.28, 18, BLACK, True)
        text(slide, sub, 8.78, y + 0.38, 2.70, 0.23, 14, MUTED)
    statement(slide, "工程纪律被保留，协作颗粒度被重塑", 0.64, 5.48, 7.15, BLUE_DARK, 20)
    tag(slide, "单位 / 讲师 / 日期待补充", 0.64, 6.37, 2.62)


def build():
    prs = Presentation()
    prs.slide_width = I(SW)
    prs.slide_height = I(SH)
    prs.core_properties.title = "传统软件开发模式与 AI Agent 时代的软件开发模式 - 新版完整母版"
    prs.core_properties.author = "Codex"
    prs.core_properties.subject = "技术管理培训"

    # P01
    cover(prs)

    # P02
    s = page(prs, "开场  /  课程结论", 2, "Agent 加快工作，", "本课程的基本立场", True)
    text(s, "人仍对结果负责", 0.62, 1.62, 7.10, 0.56, 38, BLUE, True)
    items = [("目标", "仍需人定义", BLUE), ("边界", "仍需人批准", ORANGE), ("结果", "仍需有证据", GREEN)]
    for idx, (head, detail, accent) in enumerate(items):
        x = 0.68 + idx * 3.78
        rule(s, x, 3.08, 0.60, accent, 0.05)
        text(s, head, x, 3.37, 2.70, 0.34, 22, accent, True)
        text(s, detail, x, 3.86, 2.90, 0.30, 18, TEXT)
    line(s, 0.70, 5.12, 12.13, 5.12, LINE, 1)
    tag(s, "执行可委派", 1.04, 5.62, 2.26, BLUE, BLUE_PALE)
    text(s, "+", 3.63, 5.68, 0.25, 0.24, 18, MUTED, True)
    tag(s, "验证有依据", 4.22, 5.62, 2.26, GREEN, GREEN_PALE)
    text(s, "+", 6.80, 5.68, 0.25, 0.24, 18, MUTED, True)
    tag(s, "批准有人负责", 7.41, 5.62, 2.52, ORANGE, ORANGE_PALE)

    # P03
    s = page(prs, "开场  /  课程边界", 3, "建立模式判断框架", "讨论工作组织与责任，不做工具推介")
    card(s, 0.66, 2.02, 5.80, 3.66, "本课程聚焦", ["三种模式完整流程", "角色、产物与质量责任", "Agent 的事实边界", "管理者判断清单"], BLUE, BLUE_PALE)
    card(s, 6.82, 2.02, 5.80, 3.66, "本次不展开", ["工具安装与命令", "产品优劣排名", "团队实施方案", "自主组织预测"], ORANGE, WHITE)

    # P04
    s = page(prs, "开场  /  教学案例", 4, "同一业务目标，不同开发方式", "贯穿案例：为业务系统新增客户反馈功能")
    rect(s, 0.72, 2.04, 5.18, 1.56, BLUE_PALE, None, True)
    tag(s, "用户端", 0.94, 2.28, 0.78)
    text(s, "提交反馈内容", 0.94, 2.79, 3.66, 0.37, 24, BLACK, True)
    text(s, "表达使用意见", 0.94, 3.25, 3.66, 0.24, 14, MUTED)
    rect(s, 6.36, 2.04, 5.52, 1.56, GREEN_PALE, None, True)
    tag(s, "管理员端", 6.59, 2.28, 1.02, GREEN, GREEN_PALE)
    text(s, "查看反馈列表", 6.59, 2.79, 3.80, 0.37, 24, BLACK, True)
    text(s, "了解已提交意见", 6.59, 3.25, 3.80, 0.24, 14, MUTED)
    text(s, "明确排除范围", 0.76, 4.38, 1.22, 0.23, 12, ORANGE, True)
    for idx, item in enumerate(["不含处理状态", "不含通知", "不含统计报表", "不含复杂权限"]):
        tag(s, item, 0.76 + idx * 2.82, 4.86, 2.30, ORANGE, ORANGE_PALE)
    statement(s, "范围固定，比较的才是开发方式", 0.76, 5.73, 5.80, BLUE_DARK, 19)

    # P05
    s = page(prs, "开场  /  比较坐标", 5, "我们如何比较三种模式", "七个问题贯穿完整课程")
    labels = ["需求\n表达", "工作\n组织", "参与\n角色", "交付\n产物", "质量\n验证", "发布\n责任", "反馈\n回流"]
    line(s, 0.92, 3.35, 12.14, 3.35, BLUE, 2)
    for idx, label in enumerate(labels):
        x = 0.92 + idx * 1.77
        accent = GREEN if idx == 4 else ORANGE if idx == 5 else BLUE
        dot(s, x, 3.12, 0.46, WHITE, accent)
        text(s, f"{idx + 1:02d}", x, 3.27, 0.46, 0.15, 10, accent, True, PP_ALIGN.CENTER)
        text(s, label, x - 0.36, 2.27, 1.18, 0.55, 17, BLACK, True, PP_ALIGN.CENTER)
    statement(s, "不只问是否更快，还要问如何证明正确", 0.86, 5.38, 9.60, BLUE_DARK, 22)

    # P06
    s = page(prs, "模式一  /  瀑布与 V 模型", 6, "用阶段控制复杂度", "先建立基线，再逐阶段交付")
    flow_nodes(s, ["需求确认", "设计批准", "开发实现", "测试验收", "发布交付"], 0.76, 2.28, 11.72, BLUE, 15)
    text(s, "V 型验证对应关系", 0.78, 3.38, 2.05, 0.22, 12, MUTED, True)
    line(s, 1.15, 4.09, 3.20, 5.05, BLUE, 1.5)
    line(s, 3.20, 5.05, 5.18, 4.09, BLUE, 1.5)
    text(s, "需求", 0.87, 3.74, 0.74, 0.27, 16, BLUE_DARK, True)
    text(s, "实现", 2.76, 5.18, 0.88, 0.27, 16, BLUE_DARK, True)
    text(s, "验收", 4.77, 3.74, 0.74, 0.27, 16, GREEN, True)
    card(s, 6.70, 3.50, 2.70, 1.80, "价值", ["范围可追踪", "责任易确认"], BLUE, WHITE, 15)
    card(s, 9.67, 3.50, 2.70, 1.80, "限制", ["反馈相对晚", "变更成本高"], ORANGE, WHITE, 15)

    # P07
    s = page(prs, "模式一  /  流程全景", 7, "阶段式交付全景", "案例：客户反馈功能")
    panorama(s, ["业务负责人", "需求分析", "设计开发", "测试", "发布负责"],
             ["范围确认", "需求基线", "设计评审", "编码实现", "测试验收", "发布批准"],
             ["需求规格", "设计说明", "测试用例", "验收记录", "批准记录"],
             BLUE, "只实现：提交反馈 + 查看列表")

    # P08
    s = page(prs, "模式一  /  拟真产物", 8, "产物样例：需求条目与设计说明",
             "基线先稳定，实施再展开")
    doc_sheet(s, 0.72, 2.04, 5.62, 3.72, "需求规格说明  /  v1.0", [
        ("RQ-FB-01", "用户可提交反馈内容。"),
        ("RQ-FB-02", "管理员可查看反馈列表。"),
        ("OUT-01", "不含通知与状态处理。"),
    ], BLUE)
    doc_sheet(s, 6.76, 2.04, 5.62, 3.72, "设计说明节选  /  v1.0", [
        ("UI", "反馈提交页、反馈列表页"),
        ("DATA", "反馈内容、提交时间"),
        ("RULE", "列表仅显示已提交记录"),
    ], BLUE)
    tag(s, "可评审的范围基线", 0.72, 6.17, 2.18, BLUE, BLUE_PALE)

    # P09
    s = page(prs, "模式一  /  拟真产物", 9, "产物样例：需求如何被验证",
             "V 模型强调测试与前序承诺对应")
    table(s, 0.74, 2.15, [2.10, 2.10, 4.66, 1.28], ["需求/边界", "验证编号", "检查结果", "结论"], [
        ["RQ-FB-01", "TC-FB-01", "提交后生成反馈记录", "通过"],
        ["RQ-FB-02", "TC-FB-02", "列表可看到已提交记录", "通过"],
        ["OUT-01", "RV-FB-01", "未出现通知或状态项", "通过"],
    ], 0.62, 14, GREEN)
    statement(s, "每项范围承诺，都能对应验证证据", 0.78, 5.63, 8.80, GREEN, 22)

    # P10
    s = page(prs, "模式一  /  质量责任", 10, "阶段门禁由谁把关",
             "质量不只是测试部门的工作")
    gate_chain(s, [("需求评审", "业务负责人\n确认范围", BLUE),
                   ("设计评审", "技术负责人\n批准方案", BLUE),
                   ("测试验收", "测试负责人\n提供证据", GREEN),
                   ("发布决定", "发布责任人\n批准上线", ORANGE)],
               "质量来自基线、评审与对应验证")

    # P11
    s = page(prs, "模式一  /  小结", 11, "稳定范围下的强控制方式",
             "它解决的是责任可追踪与结果可验收")
    card(s, 0.74, 2.25, 5.44, 2.70, "适合发挥价值时", ["范围相对稳定", "审查要求明确", "交付需要追踪"], BLUE, BLUE_PALE)
    card(s, 6.70, 2.25, 5.44, 2.70, "可能遇到压力时", ["用户反馈频繁", "变更不断进入", "重新确认耗时"], ORANGE, WHITE)
    statement(s, "后续模式不是丢弃基线，而是让反馈更早", 0.76, 5.60, 10.80, BLUE_DARK, 21)

    # P12
    s = page(prs, "模式二  /  敏捷与 DevOps", 12, "让反馈更早到达", "以可工作的增量验证价值")
    stages = [("计划", "故事"), ("构建", "增量"), ("测试", "证据"), ("发布", "门禁"), ("反馈", "回流")]
    for idx, (head, sub) in enumerate(stages):
        x = 0.92 + idx * 2.28
        dot(s, x, 3.06, 0.62, BLUE if idx < 3 else GREEN)
        text(s, f"{idx + 1:02d}", x, 3.27, 0.62, 0.17, 10, WHITE, True, PP_ALIGN.CENTER)
        text(s, head, x - 0.22, 2.35, 1.08, 0.28, 19, BLACK, True, PP_ALIGN.CENTER)
        text(s, sub, x - 0.30, 3.92, 1.24, 0.27, 15, MUTED, False, PP_ALIGN.CENTER)
        if idx < 4:
            line(s, x + 0.65, 3.38, x + 2.16, 3.38, BLUE, 1.5)
    line(s, 10.30, 4.60, 1.22, 4.60, GREEN, 1.2)
    text(s, "运行反馈回到下一轮计划", 4.14, 4.78, 4.78, 0.26, 16, GREEN, True, PP_ALIGN.CENTER)
    statement(s, "不是少做质量，而是更频繁验证", 0.80, 5.60, 8.20, BLUE_DARK, 21)

    # P13
    s = page(prs, "模式二  /  流程全景", 13, "迭代式交付全景", "案例：客户反馈功能")
    panorama(s, ["产品负责", "开发团队", "测试协作", "运维发布", "业务评审"],
             ["故事入列", "验收澄清", "开发测试", "流水线", "演示反馈", "发布"],
             ["用户故事", "验收标准", "任务状态", "验证记录", "反馈待办"],
             BLUE, "本轮增量：提交反馈 + 查看列表")
    line(s, 11.90, 4.95, 11.90, 5.72, GREEN, 1)
    line(s, 11.90, 5.72, 4.32, 5.72, GREEN, 1)

    # P14
    s = page(prs, "模式二  /  拟真产物", 14, "产物样例：故事、验收标准与任务板",
             "价值表达被拆成短周期可检查工作")
    doc_sheet(s, 0.72, 2.02, 5.42, 3.54, "用户故事与验收标准", [
        ("ST-01", "用户提交反馈，表达意见。"),
        ("AC-01", "提交后系统保存反馈。"),
        ("AC-02", "管理员列表可见记录。"),
    ], BLUE)
    rect(s, 6.60, 2.02, 5.78, 3.54, WHITE, LINE, True)
    text(s, "迭代任务板  /  当前增量", 6.82, 2.24, 4.20, 0.24, 13, BLUE, True)
    for idx, (status, task, accent, fill) in enumerate([
        ("TODO", "提交页体验检查", BLUE, BLUE_PALE),
        ("DOING", "列表接口实现", ORANGE, ORANGE_PALE),
        ("DONE", "验收测试", GREEN, GREEN_PALE),
    ]):
        yy = 2.77 + idx * 0.76
        tag(s, status, 6.84, yy, 0.92, accent, fill)
        text(s, task, 8.01, yy + 0.05, 3.70, 0.27, 15, TEXT, True)
    statement(s, "每个增量都以验收条件结束", 0.76, 5.91, 7.86, GREEN, 20)

    # P15
    s = page(prs, "模式二  /  拟真产物", 15, "产物样例：流水线证据与发布反馈",
             "变更更频繁，证据也更及时")
    doc_sheet(s, 0.72, 2.04, 5.48, 3.54, "Pipeline Run  #128", [
        ("BUILD", "构建通过"),
        ("TEST", "自动化测试 12 / 12 通过"),
        ("SCOPE", "范围检查通过"),
        ("RELEASE", "发布审批待确认"),
    ], GREEN)
    doc_sheet(s, 6.68, 2.04, 5.70, 3.54, "演示反馈记录", [
        ("OK", "反馈列表信息清晰"),
        ("NEXT", "统计建议进入后续讨论"),
        ("BOUND", "本轮不扩大范围"),
    ], BLUE)
    tag(s, "反馈可接收，不自动扩大本轮范围", 0.74, 5.92, 3.60, ORANGE, ORANGE_PALE)

    # P16
    s = page(prs, "模式二  /  质量责任", 16, "高频交付仍需门禁",
             "自动验证提供证据，责任人决定发布")
    gate_chain(s, [("验收标准", "产品负责\n确认价值", BLUE),
                   ("开发测试", "团队持续\n保障质量", BLUE),
                   ("流水线", "自动验证\n输出证据", GREEN),
                   ("发布授权", "责任人\n决定上线", ORANGE),
                   ("反馈排序", "产品负责\n决定后续", BLUE)],
               "频繁发布，不等于自动发布")

    # P17
    s = page(prs, "模式二  /  小结", 17, "缩短假设到反馈的距离",
             "更早看见价值，也更早看见问题")
    card(s, 0.72, 2.22, 5.46, 2.76, "获得的能力", ["更早验证价值", "更快吸收变化", "运行反馈回流"], GREEN, GREEN_PALE)
    card(s, 6.68, 2.22, 5.46, 2.76, "需要的纪律", ["持续整理待办", "自动验证可信", "发布权限清晰"], BLUE, WHITE)
    statement(s, "加速必须建立在持续验证之上", 0.74, 5.64, 8.80, BLUE_DARK, 22)

    # P18
    s = page(prs, "模式三  /  AI Agent 协作", 18, "Agent 进入真实工作上下文",
             "当前可管理、可落地的协作定义", True)
    card(s, 0.72, 2.26, 5.54, 2.66, "Agent 可执行", ["读取授权上下文", "完成明确任务", "修改产物并验证"], BLUE, BLUE_PALE)
    card(s, 6.74, 2.26, 5.54, 2.66, "人必须负责", ["设定目标边界", "审查验证证据", "批准并承担结果"], ORANGE, WHITE)
    statement(s, "可委派执行，不让渡责任", 0.74, 5.58, 7.76, BLUE_DARK, 23)

    # P19
    s = page(prs, "模式三  /  流程全景", 19, "Agent 协作全景", "案例：客户反馈功能")
    panorama(s, ["业务决策者", "技术责任人", "AI Agent", "审查发布"],
             ["确认边界", "提供上下文", "形成任务", "执行修改", "工具验证", "人工审批"],
             ["范围约束", "确认规格", "任务清单", "修改摘要", "验证结果", "批准记录"],
             GREEN, "越界扩展必须被拦截")

    # P20
    s = page(prs, "模式三  /  规格驱动", 20, "Agent 越能执行，规格越需显性",
             "速度会放大输入质量")
    card(s, 0.72, 2.09, 5.44, 3.38, "输入模糊的风险", ["模糊目标 -> 快速偏航", "隐含范围 -> 越界实现", "无验收 -> 完成错觉"], ORANGE, ORANGE_PALE)
    card(s, 6.70, 2.09, 5.44, 3.38, "可审查的控制", ["确认规格", "列明排除约束", "拆解任务与验证"], GREEN, GREEN_PALE)
    statement(s, "让意图成为可审查输入", 0.75, 5.90, 7.18, BLUE_DARK, 22)

    # P21
    s = page(prs, "模式三  /  Spec-Driven Development", 21, "从聊天意图到可执行规格",
             "基于公开实践观察到的共同趋势")
    text(s, "共同趋势", 0.76, 2.08, 1.22, 0.25, 13, BLUE, True)
    flow_nodes(s, ["意图显性化", "约束可检查", "任务可追踪", "验证有证据"], 0.76, 2.52, 11.56, BLUE, 15)
    text(s, "事实边界", 0.76, 3.74, 1.22, 0.25, 13, ORANGE, True)
    for idx, item in enumerate(["不同框架路径不同", "不是统一行业标准", "不必然替代传统文档"]):
        tag(s, item, 0.76 + idx * 3.86, 4.20, 3.38, ORANGE, ORANGE_PALE)
    statement(s, "规格更显性，执行更可控", 0.78, 5.47, 7.40, BLUE_DARK, 22)
    text(s, "来源：Spec Kit、OpenSpec、Superpowers 官方公开资料，访问于 2026-05-27",
         0.78, 6.38, 11.40, 0.18, 10, MUTED)

    # P22
    s = page(prs, "模式三  /  规格驱动", 22, "客户反馈功能：从目标到证据",
             "以确认范围为起点，以审批记录为终点")
    chain = [("目标", "收集意见", BLUE), ("范围", "提交 / 列表", BLUE), ("约束", "不含扩展", ORANGE),
             ("任务", "实现并验证", BLUE), ("证据", "测试与审查", GREEN)]
    line(s, 0.96, 3.29, 12.10, 3.29, LINE, 2)
    for idx, (head, sub, accent) in enumerate(chain):
        x = 1.02 + idx * 2.26
        dot(s, x, 3.03, 0.54, accent)
        text(s, f"{idx + 1:02d}", x, 3.22, 0.54, 0.14, 9, WHITE, True, PP_ALIGN.CENTER)
        text(s, head, x - 0.18, 2.38, 0.90, 0.28, 17, accent, True, PP_ALIGN.CENTER)
        text(s, sub, x - 0.50, 3.92, 1.60, 0.30, 14, TEXT, False, PP_ALIGN.CENTER)
    tag(s, "人确认规格", 0.82, 5.05, 1.72, ORANGE, ORANGE_PALE)
    tag(s, "Agent 执行区域", 4.74, 5.05, 2.00, BLUE, BLUE_PALE)
    tag(s, "人批准结果", 10.34, 5.05, 1.72, ORANGE, ORANGE_PALE)
    statement(s, "边界清楚，Agent 才能稳定加速", 0.82, 5.76, 9.32, BLUE_DARK, 21)

    # P23
    s = page(prs, "模式三  /  拟真产物", 23, "产物样例：先约束，再委派",
             "Agent 开始执行前的工作包")
    doc_sheet(s, 0.72, 2.00, 5.66, 3.85, "context-and-spec.md  /  已确认", [
        ("目标", "收集用户反馈供管理员查看。"),
        ("范围", "提交反馈；查看列表。"),
        ("禁止", "通知、状态、统计不在本次范围。"),
    ], BLUE)
    doc_sheet(s, 6.78, 2.00, 5.60, 3.85, "tasks.md  /  已授权", [
        ("TASK-01", "实现反馈提交能力。"),
        ("TASK-02", "实现管理员列表。"),
        ("VERIFY", "验证功能及无越界项。"),
    ], GREEN)
    tag(s, "规格已由负责人确认", 0.74, 6.16, 2.34, GREEN, GREEN_PALE)

    # P24
    s = page(prs, "模式三  /  拟真产物", 24, "产物样例：执行不能止于“完成”",
             "结果需包含修改、验证与审批记录")
    doc_sheet(s, 0.72, 2.00, 5.56, 3.94, "change-summary.md", [
        ("新增", "反馈提交表单与保存逻辑"),
        ("新增", "管理员反馈列表"),
        ("未新增", "通知、状态、统计"),
    ], BLUE)
    doc_sheet(s, 6.76, 2.00, 5.62, 3.94, "verification-and-approval.md", [
        ("功能", "提交成功；列表可见"),
        ("范围", "未发现越界项"),
        ("批准", "负责人同意本次交付"),
    ], GREEN)
    statement(s, "工具输出证据，人工确认交付", 0.76, 6.20, 8.90, GREEN, 20)

    # P25
    s = page(prs, "模式三  /  公开实践", 25, "Spec Kit：结构化规格到实现",
             "官方公开资料中的 Spec-Driven Development 实例")
    framework(s, ["constitution", "specify", "plan", "tasks", "implement"],
              ["以规格组织意图", "逐步细化任务", "实现前形成结构"],
              "来源：https://github.com/github/spec-kit  |  访问于 2026-05-27", BLUE)

    # P26
    s = page(prs, "模式三  /  公开实践", 26, "OpenSpec：为变化增加规格层",
             "轻量、迭代并面向既有系统的变更组织")
    framework(s, ["proposal", "specs", "design", "tasks"],
              ["按变更保留产物", "编码前对齐范围", "适应迭代演进"],
              "来源：https://github.com/Fission-AI/OpenSpec  |  访问于 2026-05-27", GREEN)

    # P27
    s = page(prs, "模式三  /  公开实践", 27, "Superpowers：用过程纪律约束执行",
             "设计确认、计划、测试与评审构成控制链")
    framework(s, ["澄清/设计", "计划", "执行", "测试", "评审", "完成"],
              ["先确认设计", "强调测试证据", "用评审收束结果"],
              "来源：https://github.com/obra/superpowers  |  访问于 2026-05-27", BLUE)

    # P28
    s = page(prs, "模式三  /  质量责任", 28, "Agent 需要更清晰的控制链",
             "执行可以加速，批准不能空缺")
    gate_chain(s, [("规格确认", "人确认\n目标范围", BLUE),
                   ("任务授权", "人限定\n可执行项", ORANGE),
                   ("执行修改", "Agent\n边界内处理", BLUE),
                   ("工具验证", "形成\n可核结果", GREEN),
                   ("变更审查", "人核查\n范围质量", ORANGE),
                   ("交付批准", "人承担\n最终决定", BLUE)],
               "任务可委派，责任必须归属")

    # P29
    s = page(prs, "模式三  /  失败与纠偏", 29, "纠偏案例：Agent 擅自加入通知能力",
             "完整展示错误、发现、修正、复验与批准")
    tag(s, "已确认范围", 0.78, 2.10, 1.22, BLUE, BLUE_PALE)
    text(s, "提交反馈  +  查看列表", 2.22, 2.16, 3.22, 0.22, 17, TEXT, True)
    tag(s, "错误输出", 6.62, 2.10, 1.12, ORANGE, ORANGE_PALE)
    text(s, "新增通知配置", 7.97, 2.16, 2.72, 0.22, 17, ORANGE, True)
    steps = [("发现", "识别越界", ORANGE), ("约束", "重申范围", BLUE), ("修正", "删除通知", BLUE),
             ("复验", "范围通过", GREEN), ("批准", "同意交付", GREEN)]
    line(s, 1.02, 4.10, 11.65, 4.10, LINE, 2)
    for idx, (head, body, accent) in enumerate(steps):
        x = 1.08 + idx * 2.18
        dot(s, x, 3.84, 0.52, accent)
        text(s, f"{idx + 1:02d}", x, 4.02, 0.52, 0.15, 9, WHITE, True, PP_ALIGN.CENTER)
        text(s, head, x - 0.14, 3.16, 0.82, 0.24, 16, accent, True, PP_ALIGN.CENTER)
        text(s, body, x - 0.39, 4.67, 1.34, 0.28, 14, TEXT, False, PP_ALIGN.CENTER)
    statement(s, "错误可被发现、修正并重新验证", 0.84, 5.71, 9.60, GREEN, 22)

    # P30
    s = page(prs, "横向比较  /  路径", 30, "同一案例，三种组织路径",
             "目标相同，反馈颗粒度不同")
    paths = [
        ("瀑布 / V", ["基线", "设计", "实现", "测试", "发布"], BLUE),
        ("敏捷 / DevOps", ["故事", "增量", "流水线", "反馈", "发布"], BLUE),
        ("AI Agent", ["规格/约束", "委派执行", "工具验证", "人审", "批准"], GREEN),
    ]
    for idx, (name, nodes, accent) in enumerate(paths):
        yy = 2.16 + idx * 1.22
        text(s, name, 0.76, yy + 0.15, 1.72, 0.28, 16, accent, True)
        flow_nodes(s, nodes, 2.72, yy, 9.56, accent, 13)
    statement(s, "比较的是控制结构，不是简单排名", 0.80, 5.91, 9.20, BLUE_DARK, 21)

    # P31
    s = page(prs, "横向比较  /  角色", 31, "执行方式变化，责任仍需归属",
             "谁做什么，谁负责什么")
    table(s, 0.72, 2.04, [2.20, 2.98, 3.00, 3.03], ["判断问题", "瀑布 / V", "敏捷 / DevOps", "AI Agent"], [
        ["范围决定", "业务负责人", "产品负责人", "人确认范围"],
        ["任务执行", "环节分工", "团队协作", "人与 Agent"],
        ["验证证据", "测试记录", "流水线与团队", "工具结果与人审"],
        ["发布批准", "责任人", "责任人", "责任人"],
    ], 0.62, 13, BLUE_DARK)
    statement(s, "Agent 改变执行分配，不接管最终责任", 0.76, 5.72, 10.60, BLUE_DARK, 21)

    # P32
    s = page(prs, "横向比较  /  产物", 32, "交付产物如何演进",
             "文档形态变化，可追踪要求仍在")
    groups = [
        ("瀑布 / V", ["需求规格", "设计说明", "测试追踪", "验收记录"], BLUE),
        ("敏捷 / DevOps", ["用户故事", "验收标准", "流水线证据", "反馈待办"], BLUE),
        ("AI Agent", ["确认规格", "上下文约束", "修改与验证", "人工批准"], GREEN),
    ]
    for idx, (head, values, accent) in enumerate(groups):
        x = 0.74 + idx * 4.04
        text(s, head, x, 2.06, 3.52, 0.31, 19, accent, True)
        rule(s, x, 2.53, 0.56, accent, 0.045)
        for item_idx, item in enumerate(values):
            rect(s, x, 2.84 + item_idx * 0.61, 3.45, 0.46, BLUE_PALE if accent == BLUE else GREEN_PALE,
                 None, True)
            text(s, item, x + 0.14, 2.95 + item_idx * 0.61, 3.16, 0.23, 15, TEXT, True)
    rect(s, 0.74, 5.72, 11.56, 0.56, BLUE_DARK, None, True)
    text(s, "共同基座：可审查、可验证、可承担责任的交付证据",
         0.82, 5.89, 11.40, 0.23, 18, WHITE, True, PP_ALIGN.CENTER)

    # P33
    s = page(prs, "横向比较  /  质量责任", 33, "三种模式都必须回答质量问题",
             "门禁位置改变，责任没有消失")
    table(s, 0.72, 2.06, [3.02, 2.79, 2.79, 2.79], ["必须回答的问题", "瀑布 / V", "敏捷 / DevOps", "AI Agent"], [
        ["范围如何确认", "需求评审", "验收澄清", "规格批准"],
        ["结果如何验证", "对应测试", "持续流水线", "工具验证+人审"],
        ["发布谁批准", "发布责任人", "授权责任人", "人工审批者"],
        ["反馈如何回流", "变更管理", "待办排序", "新规格/任务"],
    ], 0.60, 13, GREEN)
    tag(s, "验证频率提高", 0.82, 5.62, 1.82, GREEN, GREEN_PALE)
    text(s, "低频阶段门禁", 3.10, 5.70, 2.10, 0.23, 14, MUTED)
    line(s, 5.16, 5.82, 8.28, 5.82, BLUE, 1.4)
    text(s, "高频执行反馈", 8.64, 5.70, 2.18, 0.23, 14, GREEN, True)

    # P34
    s = page(prs, "横向比较  /  适用条件", 34, "不是三选一，而是组合能力",
             "选择协作方式，要看任务与控制能力")
    conditions = [
        ("范围稳定、审查严格", "强化基线与评审", BLUE),
        ("需要快速验证价值", "强化迭代与反馈", BLUE),
        ("上下文清晰、验证可执行", "可引入 Agent 协作", GREEN),
    ]
    for idx, (condition, response, accent) in enumerate(conditions):
        x = 0.76 + idx * 3.96
        rect(s, x, 2.24, 3.50, 2.32, WHITE, LINE, True)
        tag(s, f"条件 {idx + 1:02d}", x + 0.20, 2.51, 0.91, accent,
            GREEN_PALE if accent == GREEN else BLUE_PALE)
        text(s, condition, x + 0.20, 3.10, 3.05, 0.52, 17, BLACK, True)
        text(s, response, x + 0.20, 3.96, 3.04, 0.25, 15, accent, True)
    rect(s, 0.78, 5.35, 11.52, 0.70, BLUE_DARK, None, True)
    text(s, "真实团队通常组合这些能力，形成可靠交付", 0.84, 5.58, 11.38, 0.25, 20, WHITE, True, PP_ALIGN.CENTER)

    # P35
    s = page(prs, "管理收束  /  价值与风险", 35, "更短反馈距离，也可能更快放大错误",
             "提效能力必须配合控制能力", True)
    card(s, 0.72, 2.35, 5.46, 2.72, "潜在价值", ["快速阅读上下文", "缩短修改循环", "及时形成证据"], GREEN, GREEN_PALE)
    card(s, 6.68, 2.35, 5.46, 2.72, "必须防范", ["误解快速实现", "越界扩大范围", "无验证却完成"], ORANGE, ORANGE_PALE)
    statement(s, "真正要判断的是：团队能否控制加速", 0.76, 5.70, 10.50, BLUE_DARK, 22)

    # P36
    s = page(prs, "管理收束  /  检查清单", 36, "引入 Agent 前，先检查任务条件",
             "任何关键条件不清，先完善治理输入")
    checks = ["目标能否明确表述？", "范围能否写成约束？", "上下文是否可提供？",
              "任务能否分段验收？", "敏感权限是否可隔离？"]
    for idx, item in enumerate(checks):
        y = 2.06 + idx * 0.69
        accent = ORANGE if idx == 4 else BLUE
        dot(s, 0.88, y + 0.06, 0.30, WHITE, accent)
        text(s, "?", 0.88, y + 0.13, 0.30, 0.14, 11, accent, True, PP_ALIGN.CENTER)
        text(s, item, 1.50, y, 5.62, 0.34, 20, BLACK, True)
        line(s, 7.84, y + 0.23, 11.78, y + 0.23, LINE, 1)
        text(s, "是 / 否 / 待完善", 9.40, y, 2.18, 0.30, 13, MUTED, False, PP_ALIGN.RIGHT)

    # P37
    s = page(prs, "管理收束  /  检查清单", 37, "引入 Agent 后，谁来确认结果",
             "没有责任归属，就没有可靠委派")
    responsibilities = [
        ("规格批准", "谁确认目标与范围？", BLUE),
        ("证据复核", "谁核对验证结果？", GREEN),
        ("越界拦截", "谁处理范围偏差？", ORANGE),
        ("发布决定", "谁承担交付责任？", BLUE),
        ("记录留存", "谁保存过程证据？", GREEN),
    ]
    for idx, (head, question, accent) in enumerate(responsibilities):
        y = 2.05 + idx * 0.71
        tag(s, head, 0.88, y, 1.32, accent, GREEN_PALE if accent == GREEN else ORANGE_PALE if accent == ORANGE else BLUE_PALE)
        text(s, question, 2.66, y + 0.04, 4.42, 0.27, 18, TEXT, True)
        rect(s, 8.30, y + 0.03, 3.42, 0.29, WHITE, LINE, True)
        text(s, "责任人 / 机制：____________", 8.46, y + 0.08, 3.06, 0.18, 11, MUTED)

    # P38
    s = page(prs, "管理收束  /  课程总结", 38, "从阶段控制，到持续反馈，", None, True)
    text(s, "再到受控执行", 0.62, 1.52, 6.90, 0.54, 38, BLUE, True)
    line(s, 0.86, 3.32, 12.02, 3.32, BLUE, 2.2)
    close = [("瀑布 / V", "建立基线与追踪", BLUE), ("敏捷 / DevOps", "缩短价值反馈", BLUE),
             ("AI Agent", "提高受控执行密度", GREEN)]
    for idx, (name, message, accent) in enumerate(close):
        x = 1.00 + idx * 3.88
        dot(s, x, 3.09, 0.48, WHITE, accent)
        text(s, name, x, 2.44, 3.18, 0.31, 19, accent, True)
        text(s, message, x, 3.89, 3.28, 0.30, 16, TEXT)
    rect(s, 0.84, 5.32, 11.58, 0.78, BLUE_DARK, None, True)
    text(s, "Agent 加速软件工作，人定义方向并对结果负责", 0.98, 5.57, 11.28, 0.29, 22, WHITE, True, PP_ALIGN.CENTER)

    # P39
    s = page(prs, "参考资料  /  事实来源", 39, "参考资料", "公开实践仅用于说明取向，不构成选型推荐")
    refs = [
        ("Spec Kit", "规格驱动开发工具包", "github.com/github/spec-kit", BLUE),
        ("OpenSpec", "轻量规格与变更组织", "github.com/Fission-AI/OpenSpec", GREEN),
        ("Superpowers", "Agent 开发过程方法", "github.com/obra/superpowers", BLUE),
    ]
    for idx, (name, purpose, url, accent) in enumerate(refs):
        y = 2.15 + idx * 1.05
        rule(s, 0.82, y + 0.16, 0.58, accent, 0.045)
        text(s, name, 1.64, y, 2.38, 0.33, 20, BLACK, True)
        text(s, purpose, 4.28, y + 0.04, 2.86, 0.27, 16, TEXT)
        text(s, url, 8.08, y + 0.04, 3.98, 0.27, 14, accent, True)
    line(s, 0.82, 5.65, 12.05, 5.65, LINE, 1)
    text(s, "资料依据其官方公开仓库整理，访问日期：2026-05-27",
         0.82, 5.98, 8.20, 0.25, 14, MUTED)
    tag(s, "谢谢", 11.10, 5.92, 0.92, BLUE, BLUE_PALE)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(str(OUT))
    print(f"slides={len(prs.slides)}")
    return OUT


if __name__ == "__main__":
    build()

